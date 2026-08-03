from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from build_thesis_docx import ensure_image_asset, png_dimensions


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"
PAPERS = OUTPUTS / "papers"
DATA = OUTPUTS / "data"
FIGURES = OUTPUTS / "figures"
ARCHITECTURE = FIGURES / "architecture"
PUBLIC_BENCHMARK_FIGURES = FIGURES / "public_benchmarks"
THESIS_PAPERS = ROOT / "docs" / "papers" / "thesis"
PRESENTATION_PATH = PAPERS / "thesis_presentation_zh.pptx"
STORED_PRESENTATION_PATH = THESIS_PAPERS / "thesis_presentation_zh.pptx"
OUTLINE_PATH = ROOT / "docs" / "thesis" / "presentation_outline_zh.md"
LONG_PRESENTATION_PATH = PAPERS / "thesis_presentation_zh_30min.pptx"
STORED_LONG_PRESENTATION_PATH = THESIS_PAPERS / "thesis_presentation_zh_30min.pptx"
LONG_OUTLINE_PATH = ROOT / "docs" / "thesis" / "presentation_outline_zh_30min.md"
LONG_SPEAKER_NOTES_PATH = ROOT / "docs" / "thesis" / "presentation_speaker_notes_zh_30min.md"

BACKGROUND_COLOR = RGBColor(244, 247, 251)
HEADER_FILL = RGBColor(18, 32, 51)
HEADER_TEXT = RGBColor(255, 255, 255)
HEADER_SUBTITLE = RGBColor(211, 225, 241)
TITLE_COLOR = HEADER_TEXT
TEXT_COLOR = RGBColor(30, 39, 51)
ACCENT_COLOR = RGBColor(0, 97, 148)
ACCENT_LIGHT = RGBColor(45, 155, 200)
MUTED_COLOR = RGBColor(70, 80, 92)
CARD_FILL = RGBColor(255, 255, 255)
CARD_LINE = RGBColor(209, 219, 231)
CARD_SHADOW = RGBColor(226, 232, 240)
BODY_FONT = "Noto Sans TC"
LATIN_FONT = "Arial"
FORMULA_FONT = "Cambria Math"


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def window_temperature_domain_counts(summary: dict) -> Tuple[int, int]:
    scenarios = summary.get("scenarios", [])
    in_domain = sum(
        20.0 <= float(item["target_zone_estimated"]["temperature"]) <= 30.0
        for item in scenarios
    )
    return in_domain, len(scenarios) - in_domain


def average_field_mae(summary: dict) -> dict:
    scenarios = summary.get("scenarios", [])
    metrics = ("temperature", "humidity", "illuminance")
    return {
        metric: round(sum(item["field_mae"][metric] for item in scenarios) / max(len(scenarios), 1), 4)
        for metric in metrics
    }


def best_recommendations(summary: dict) -> List[Tuple[str, str]]:
    output = []
    for item in summary.get("scenarios", []):
        recommendations = item.get("recommendations", [])
        best = recommendations[0]["name"] if recommendations else "n/a"
        output.append((item["name"], best))
    return output[:5]


def scenario_map(summary: dict) -> dict:
    return {item["name"]: item for item in summary.get("scenarios", [])}


def init_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def style_slide(slide) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BACKGROUND_COLOR


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    return slide


def set_frame_margins(frame, left: float, right: float, top: float, bottom: float) -> None:
    frame.margin_left = Inches(left)
    frame.margin_right = Inches(right)
    frame.margin_top = Inches(top)
    frame.margin_bottom = Inches(bottom)


def fit_text_size(lines: Sequence[str], width: float, height: float, preferred: int, minimum: int = 10) -> int:
    if not lines:
        return preferred
    longest = max(len(line) for line in lines)
    wrapped_units = 0
    for line in lines:
        # Chinese glyphs are wider than Latin words in this deck; this keeps dense cards from overflowing.
        chars_per_line = max(8, int(width * 5.8 * (13 / max(preferred, 1))))
        wrapped_units += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    line_capacity = max(1, int(height * 72 / (preferred * 1.42)))
    size = preferred
    if wrapped_units > line_capacity:
        size -= min(3, wrapped_units - line_capacity)
    if longest > width * 9.0:
        size -= 1
    return max(minimum, min(preferred, size))


def bullet_text(text: str) -> str:
    if re.match(r"^\d+[.．、]\s*", text):
        return text
    return f"• {text}"


def add_title(slide, text: str, subtitle: str = "") -> None:
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.98))
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_FILL
    header.line.color.rgb = HEADER_FILL
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.95), Inches(13.333), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_LIGHT
    accent.line.color.rgb = ACCENT_LIGHT
    title_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.12), Inches(12.2), Inches(0.48))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_frame_margins(frame, 0, 0, 0, 0)
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(20 if len(text) > 24 else 22)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.26))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        sub_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        set_frame_margins(sub_frame, 0, 0, 0, 0)
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = BODY_FONT
        sub_run.font.size = Pt(10)
        sub_run.font.color.rgb = HEADER_SUBTITLE


def add_footer(slide, page: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.25), Inches(7.0), Inches(0.6), Inches(0.25))
    box.name = "footer_page_number"
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page)
    run.font.name = LATIN_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED_COLOR


def renumber_footers(prs: Presentation) -> None:
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name != "footer_page_number" or not shape.has_text_frame:
                continue
            frame = shape.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.RIGHT
            run = paragraph.add_run()
            run.text = str(index)
            run.font.name = LATIN_FONT
            run.font.size = Pt(9)
            run.font.color.rgb = MUTED_COLOR
            break


def add_styled_run(
    paragraph,
    text: str,
    size: int,
    color: RGBColor,
    bold: bool = False,
    font_name: str = BODY_FONT,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_markup_runs(
    paragraph,
    text: str,
    size: int,
    color: RGBColor,
    font_name: str = BODY_FONT,
    render_simple_subscripts: bool = False,
) -> None:
    """Render subscript markup like t_{ref}, and optionally N_T, in PowerPoint text."""
    cursor = 0
    pattern = r"_\{([^}]+)\}|_([A-Za-z0-9]+(?:,[A-Za-z0-9]+)?)" if render_simple_subscripts else r"_\{([^}]+)\}"
    for match in re.finditer(pattern, text):
        if match.start() > cursor:
            add_styled_run(paragraph, text[cursor : match.start()], size, color, font_name=font_name)
        subscript_text = match.group(1) if match.group(1) is not None else match.group(2)
        run = paragraph.add_run()
        run.text = subscript_text
        run.font.name = font_name
        run.font.size = Pt(max(size - 4, 8))
        run.font.color.rgb = color
        run._r.get_or_add_rPr().set("baseline", "-25000")
        cursor = match.end()
    if cursor < len(text):
        add_styled_run(paragraph, text[cursor:], size, color, font_name=font_name)


def metric_triplet(values: dict, decimals: int = 4) -> str:
    return (
        f"T={values['temperature']:.{decimals}f}, "
        f"H={values['humidity']:.{decimals}f}, "
        f"L={values['illuminance']:.{decimals}f}"
    )


def percent_reduction(before: float, after: float) -> float:
    if before == 0:
        return 0.0
    return (before - after) / before * 100


def is_formula_line(text: str) -> bool:
    return text.startswith(
        (
            "T(",
            "H(",
            "L(",
            "p =",
            "v ∈",
            "Fᵥ",
            "F̂",
            "b₀",
            "T₀",
            "H₀",
            "L₀",
            "ζ",
            "A_ac",
            "A_win",
            "A_light",
            "N_T",
            "N_H",
            "N_L",
            "B_T",
            "S_T",
            "B_H",
            "S_H",
            "B_ac",
            "S_ac",
            "B_win",
            "S_win",
            "B_light",
            "S_light",
            "L_win",
            "L_light",
            "A_j",
            "Aⱼ",
            "E_j",
            "Eⱼ",
            "R_j",
            "Rⱼ",
            "Cᵥ",
            "C_v",
            "C(p)",
            "rᵛ",
            "R_v",
            "Rᵥ",
            "Δy",
            "X_",
            "β",
            "q_m",
            "P(q",
            "m ∈",
            "Corr",
            "Penalty",
            "F_hybrid",
            "Fᵛ",
            "F_true",
            "Loss",
            "ℒ",
            "φᵢ",
            "yᵢ",
            "I^",
            "Iʳ",
            "𝒱",
            "|Rᵥ",
            "MAE",
            "RMSE",
            "IDW",
            "w_s",
            "Score",
            "×",
            "+",
            "-",
            "+ Σ",
        )
    )


def add_bullets(slide, left: float, top: float, width: float, height: float, items: Sequence[str], level0_size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_frame_margins(frame, 0.02, 0.02, 0.02, 0.02)
    effective_size = fit_text_size(items, width, height, level0_size, minimum=12)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(max(4, effective_size * 0.45))
        p.line_spacing = 1.12
        p.font.name = BODY_FONT
        p.font.size = Pt(effective_size)
        p.font.color.rgb = TEXT_COLOR
        add_styled_run(p, bullet_text(item), effective_size, TEXT_COLOR)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body_lines: Sequence[str],
    title_size: int = 16,
    body_size: int = 12,
    formula_size: int = 14,
    render_simple_subscripts: bool = False,
) -> None:
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.03),
        Inches(top + 0.04),
        Inches(width),
        Inches(height),
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = CARD_SHADOW
    shadow.line.color.rgb = CARD_SHADOW
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(0.8)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_LIGHT
    accent.line.color.rgb = ACCENT_LIGHT

    title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.18), Inches(width - 0.4), Inches(0.36))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_frame_margins(title_frame, 0, 0, 0, 0)
    title_para = title_frame.paragraphs[0]
    title_para.font.name = BODY_FONT
    title_para.font.size = Pt(title_size)
    title_para.font.bold = True
    title_para.font.color.rgb = ACCENT_COLOR
    add_styled_run(title_para, title, title_size, ACCENT_COLOR, bold=True)

    body_top = top + 0.66
    body_height = max(0.2, height - 0.84)
    body_width = max(0.4, width - 0.42)
    body_box = slide.shapes.add_textbox(Inches(left + 0.21), Inches(body_top), Inches(body_width), Inches(body_height))
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_frame_margins(frame, 0, 0, 0, 0)
    non_formula_lines = [line for line in body_lines if not is_formula_line(line)]
    effective_body_size = fit_text_size(non_formula_lines or body_lines, body_width, body_height, body_size, minimum=10)
    for line in body_lines:
        para = frame.paragraphs[0] if frame.paragraphs[0].text == "" else frame.add_paragraph()
        is_formula = is_formula_line(line)
        para.font.name = FORMULA_FONT if is_formula else BODY_FONT
        para.font.size = Pt(formula_size if is_formula else effective_body_size)
        para.font.color.rgb = TEXT_COLOR
        para.space_after = Pt(2 if is_formula else 4)
        para.line_spacing = 1.08 if is_formula else 1.14
        add_markup_runs(
            para,
            line if is_formula else bullet_text(line),
            formula_size if is_formula else effective_body_size,
            TEXT_COLOR,
            font_name=FORMULA_FONT if is_formula else BODY_FONT,
            render_simple_subscripts=render_simple_subscripts,
        )


def add_learning_record_slide(prs: Presentation, footer_number: int) -> None:
    slide = new_slide(prs)
    add_title(slide, "learn_impacts：動作如何成為資料記錄")
    add_card(
        slide,
        0.7,
        1.35,
        3.85,
        5.0,
        "1. start：操作狀態",
        [
            "device_name：ac_main / window_main / light_main",
            "device_state：activation、kind、power",
            "AC state：mode、setpoint、fan、風向、swing",
            "合併成 device_specs 並更新 runtime state",
        ],
        body_size=12,
    )
    add_card(
        slide,
        4.75,
        1.35,
        3.85,
        5.0,
        "2. record：情境快照",
        [
            "learning_record_id + status=RECORDING",
            "baseline、outdoor boundary、elapsed time",
            "furniture / obstruction state",
            "before_observations：sensor -> T/H/L",
            "optional note 與 sample point 預測",
        ],
        body_size=12,
    )
    add_card(
        slide,
        8.8,
        1.35,
        3.8,
        5.0,
        "3. finish：轉成係數",
        [
            "after_observations：同一批 sensors",
            "Delta y = after - before",
            "influence envelope 作為 X",
            "least squares 解 metric coefficients",
            "輸出 learned_device_impacts 與 sensor MAE",
        ],
        body_size=12,
    )
    add_footer(slide, footer_number)


def add_picture(slide, source: Path, left: float, top: float, width: float, height: float) -> None:
    png = ensure_image_asset({"path": str(source.relative_to(ROOT)), "asset_name": source.stem})
    image_width_px, image_height_px = png_dimensions(png)
    image_ratio = image_width_px / image_height_px
    box_ratio = width / height
    if image_ratio >= box_ratio:
        fitted_width = width
        fitted_height = width / image_ratio
    else:
        fitted_height = height
        fitted_width = height * image_ratio
    fitted_left = left + (width - fitted_width) / 2
    fitted_top = top + (height - fitted_height) / 2
    slide.shapes.add_picture(
        str(png),
        Inches(fitted_left),
        Inches(fitted_top),
        width=Inches(fitted_width),
        height=Inches(fitted_height),
    )


def add_two_column_title_body(slide, title: str, left_items: Sequence[str], right_image: Path, subtitle: str = "") -> None:
    add_title(slide, title, subtitle)
    add_bullets(slide, 0.8, 1.55, 5.0, 4.9, left_items, level0_size=18)
    add_picture(slide, right_image, 6.3, 1.45, 6.2, 4.9)


FORMULA_WALKTHROUGH = [
    (
        "公式說明 1：三因子場與查詢點",
        "場的定義",
        [
            "T(p,t)：位置 p、時間 t 的溫度",
            "H(p,t)：位置 p、時間 t 的相對濕度",
            "L(p,t)：位置 p、時間 t 的照度",
            "p = (x,y,z)，單位為公尺",
            "t 可代表啟動後時間、情境時間或 demo 時間軸",
        ],
        "適用範圍",
        [
            "本研究不是只估單一平均值",
            "輸出是三維空間中任意點的三個環境量",
            "8 顆角落感測器只提供稀疏觀測",
            "其他採樣點是由模型與校正場推估出來",
            "其他採樣點為模型估計值，非直接量測值",
        ],
    ),
    (
        "公式說明 2：總估計式",
        "主公式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "v ∈ {T,H,L}",
            "T：temperature；H：relative humidity；L：illuminance",
            "Nᵥ：nominal model，描述主要物理趨勢",
            "Cᵥ：8 顆角落感測器 residual 形成的校正場",
        ],
        "為什麼這樣拆",
        [
            "Nᵥ 讓模型先有設備、邊界與空間結構",
            "Cᵥ 讓模型對齊真實角落感測器",
            "三個變數共用此估計框架",
            "但 N_T、N_H、N_L 的物理項目分開設計",
            "這能回應「不能把同一套 bulk/local 硬套三種物理量」",
        ],
    ),
    (
        "公式說明 3：Indoor baseline",
        "baseline 定義",
        [
            "b₀ = (T₀, H₀, L₀)",
            "T₀：設備作用前的室內基準溫度",
            "H₀：設備作用前的室內基準相對濕度",
            "L₀：設備作用前的室內基準照度",
            "若有啟動前觀測，可由 8 顆感測器平均取得",
        ],
        "跟 baseline 比較法的差別",
        [
            "這裡的 baseline 是模型起始狀態",
            "不是第 5 章的 IDW baseline",
            "也不是公開資料集的 persistence 或 linear regression",
            "Web demo 左側 Indoor Baseline 就是在設定 T₀/H₀/L₀",
            "所有設備影響都是在 b₀ 上加減偏移",
        ],
    ),
    (
        "公式說明 4：baseline 的取得方式",
        "有啟動前觀測時",
        [
            "T₀ = (1/|S|) ∑ₛ∈S O_{T}(pₛ, t_{ref})",
            "H₀ = (1/|S|) ∑ₛ∈S O_{H}(pₛ, t_{ref})",
            "L₀ = (1/|S|) ∑ₛ∈S O_{L}(pₛ, t_{ref})",
            "S 是 8 顆角落感測器集合",
            "t_{ref} 是設備尚未加入作用的參考時間",
        ],
        "沒有啟動前觀測時",
        [
            "改由房間設計檔、情境設定或 demo 輸入提供",
            "例如標準房間預設 T₀=29°C、H₀=67%、L₀=90 lux",
            "因此 baseline 不是模型學出來的黑盒值",
            "它是後續設備影響與 residual correction 的共同起點",
            "baseline 資料來源需在情境設定中明確標示",
        ],
    ),
    (
        "公式說明 5：高度正規化",
        "垂直座標",
        [
            "ζ = z / Hᵣ - 1/2",
            "Hᵣ：房間高度",
            "z：查詢點高度",
            "ζ 約落在 -0.5 到 0.5",
            "ζ > 0 表示偏上方；ζ < 0 表示偏下方",
        ],
        "為什麼需要",
        [
            "室內溫度與濕度可能存在垂直分層",
            "冷空氣、熱源與混合程度會讓上下層不同",
            "ζ 提供低成本的高度修正項",
            "照度主要由光源幾何與遮蔽處理",
            "所以高度不是三種變數完全相同地使用",
        ],
    ),
    (
        "公式說明 6：設備 activation",
        "時間響應",
        [
            "Aⱼ(t) = aⱼ(1 - exp(-t/τⱼ))",
            "j 代表某個設備，例如冷氣、窗戶或燈具",
            "aⱼ：設備影響的穩態比例或強度尺度",
            "τⱼ：接近穩態所需的時間常數",
            "t 越大，Aⱼ(t) 越接近 aⱼ",
        ],
        "使用原因",
        [
            "設備不會在啟動瞬間把全室改變到穩態",
            "冷氣降溫、除濕與窗戶交換都需要時間",
            "這是一階收斂近似，計算成本低且可解釋",
            "不是完整 HVAC transient simulation",
            "目標是表達主要時間趨勢",
        ],
    ),
    (
        "公式說明 7：influence envelope",
        "空間作用範圍",
        [
            "Eⱼ(p,t) = Aⱼ(t) Rⱼ(p) Dⱼ(p,t) Vⱼ(p)",
            "Aⱼ(t)：設備目前啟動強度",
            "Rⱼ(p)：距離衰減",
            "Dⱼ(p,t)：方向性，例如冷氣出風或窗戶日照方向",
            "Vⱼ(p)：可見性或遮蔽程度",
        ],
        "距離衰減",
        [
            "Rⱼ(p) = exp(-||p - pⱼ|| / rⱼ)",
            "pⱼ：設備位置",
            "rⱼ：設備作用半徑或衰減尺度",
            "距離越遠，局部影響越小",
            "但全室平均項仍由各變數自己的 B 項處理",
        ],
    ),
    (
        "公式說明 8：溫度場主式",
        "溫度 nominal model",
        [
            "N_T(p,t) = T₀ + B_T(t) + S_T(p,t) + γ_T M(t) ζ",
            "T₀：室內基準溫度",
            "B_T(t)：全室平均熱響應",
            "S_T(p,t)：局部空間熱響應",
            "γ_T M(t)ζ：垂直溫度分層",
        ],
        "使用原因",
        [
            "溫度受熱交換、熱源與空氣混合影響",
            "只做局部衰減會低估冷氣的全室降溫",
            "所以保留 B_T 表示整體室溫移動",
            "S_T 再表達出風口、窗邊或燈具附近差異",
            "這是控制導向 reduced-order 熱場近似",
        ],
    ),
    (
        "公式說明 9：溫度的全室與局部項",
        "分解式",
        [
            "B_T(t) = B_ac,T(t) + B_win,T(t) + B_light,T(t)",
            "S_T(p,t) = S_ac,T(p,t) + S_win,T(p,t) + S_light,T(p,t)",
            "B_T 負責全室平均狀態改變",
            "S_T 負責某些位置比較強的局部差異",
            "兩者合起來避免 local-only 或 well-mixed-only 的偏誤",
        ],
        "三類來源",
        [
            "冷氣：依模式與設定溫差讓室內趨冷或趨暖",
            "窗戶：依 T_out - T₀ 表示外氣熱交換方向",
            "燈具：在溫度路徑中視為小型熱源",
            "燈具在溫度模型中代表發熱",
            "這就是變數專屬公式的意義",
        ],
    ),
    (
        "公式說明 10：冷氣溫度項",
        "冷氣全室項",
        [
            "B_ac,T(t) = s_m k_ac,Tᵍ d_T P_ac A_ac(t)",
            "s_m：冷房或暖房模式符號",
            "k_ac,Tᵍ：冷氣對全室溫度的增益係數",
            "d_T：設定溫度與室內基準的需求差",
            "P_ac：校正後冷氣 power scale",
        ],
        "冷氣局部項",
        [
            "S_ac,T(p,t) = s_m k_ac,Tˢ d_T P_ac E_ac(p,t)",
            "k_ac,Tˢ：冷氣局部空間增益",
            "E_ac(p,t)：出風口附近、方向與遮蔽造成的空間權重",
            "B 表示全室趨勢，S 表示出風口附近差異",
            "兩者不是任意疊加，而是修正 local-only 的缺陷",
        ],
    ),
    (
        "公式說明 11：窗戶與燈具溫度項",
        "窗戶熱交換",
        [
            "B_win,T(t) = k_win,Tᵍ (T_out - T₀) P_win A_win(t)",
            "S_win,T(p,t) = k_win,Tˢ (T_out - T₀) P_win E_win(p,t)",
            "T_out > T₀ 時偏升溫",
            "T_out < T₀ 時偏降溫",
            "P_win 表示窗戶開啟比例或校正後影響尺度",
        ],
        "燈具熱源",
        [
            "B_light,T(t) = k_light,Tᵍ P_light A_light(t)",
            "S_light,T(p,t) = k_light,Tˢ P_light E_light(p,t)",
            "燈具在溫度模型裡只代表發熱",
            "照明造成的 lux 變化由照度模型另外處理",
            "這可避免把光學效果誤解成熱場效果",
        ],
    ),
    (
        "公式說明 12：濕度場主式",
        "濕度 nominal model",
        [
            "N_H(p,t) = clip[0,100]{H₀ + B_H(t) + S_H(p,t) - γ_H M(t) ζ}",
            "H₀：室內基準相對濕度",
            "B_H(t)：全室平均水氣或除濕響應",
            "S_H(p,t)：局部水氣交換或除濕差異",
            "clip[0,100]：相對濕度限制在 0% 到 100%",
        ],
        "使用原因",
        [
            "相對濕度有明確物理範圍",
            "冷氣常見效果是除濕，所以符號方向不同於溫度",
            "窗戶則由室外濕度與室內基準濕度差決定",
            "本研究不主張完整 psychrometric model",
            "而是用低階水氣交換近似再由 sensor residual 校正",
        ],
    ),
    (
        "公式說明 13：濕度來源項",
        "全室濕度項",
        [
            "B_H(t) = -k_ac,Hᵍ d_H P_ac A_ac(t)",
            "+ k_win,Hᵍ (H_out - H₀) P_win A_win(t)",
            "冷氣項為負：表示除濕",
            "窗戶項正負由 H_out - H₀ 決定",
            "外面較濕時開窗提高濕度，較乾時降低濕度",
        ],
        "局部濕度項",
        [
            "S_H(p,t) = -k_ac,Hˢ d_H P_ac E_ac(p,t)",
            "+ k_win,Hˢ (H_out - H₀) P_win E_win(p,t)",
            "E_ac 讓除濕效果在冷氣影響區附近更強",
            "E_win 讓窗邊水氣交換較強",
            "濕度沒有使用燈具照度那套光學公式",
        ],
    ),
    (
        "公式說明 14：照度場主式",
        "照度 nominal model",
        [
            "N_L(p,t) = max{0, L₀ + L_winᵈⁱʳ(p,t)",
            "+ L_lightᵈⁱʳ(p,t) + L_winᵃᵐᵇ(p,t) + Iʳᵉᶠˡ(p,t)}",
            "L₀：室內基準照度",
            "max{0,...}：照度不可為負",
            "照度由燈具光束、窗戶、遮蔽與反射決定",
        ],
        "為什麼不同於溫濕度",
        [
            "照度不是空氣混合或水氣交換問題",
            "它更接近光線幾何與可視性問題",
            "燈具與窗戶可造成局部高照度峰值",
            "所以保留 direct source 與 obstruction",
            "再用一次漫反射補足間接光",
        ],
    ),
    (
        "公式說明 15：直射光與環境光",
        "窗戶直射光",
        [
            "L_winᵈⁱʳ(p,t) = S_out d_f k_sol P_win E_win(p,t)",
            "S_out：室外日照強度",
            "d_f：與時間、季節或日照方向相關的折減",
            "k_sol：窗戶日照轉換係數",
            "E_win：窗戶到室內點的距離、方向與遮蔽權重",
        ],
        "燈具與環境光",
        [
            "L_lightᵈⁱʳ(p,t) = G_light P_light A_light(t) Φ_light(p) Q_light(p) V_light(p)",
            "Φ_light：由光束角推得的 cosine 方向權重",
            "Q_light：參考距離正規化後的距離衰減；V_light：遮蔽或可見性",
            "L_winᵃᵐᵇ(p,t)：窗戶帶來的擴散環境光",
            "直射與環境光分開，可描述窗邊與全室背景亮度",
        ],
    ),
    (
        "公式說明 16：一次漫反射",
        "反射公式",
        [
            "Iʳᵉᶠˡ(p,t) = Σ_s ρ_s Ī_s A_sʳᵉˡ exp(-||p-c_s||/ℓ_s)",
            "× max(0, n_s·r̂_s→p) V_s(p)",
            "s：牆、地板、天花板或家具表面",
            "ρ_s：表面反射率",
            "Ī_s：表面接收到的平均照度",
        ],
        "模型限制",
        [
            "一次漫反射用來補足非直射區域的回填亮度",
            "它不是完整 ray tracing 或 radiosity",
            "只計算一次反射，因此成本較低",
            "主要用途是改善照度空間趨勢",
            "不等同精密光學模擬等級",
        ],
    ),
    (
        "公式說明 17：8 參數校正多項式",
        "三線性形式",
        [
            "C(p) = c₀ + c₁X + c₂Y + c₃Z",
            "+ c₄XY + c₅XZ + c₆YZ + c₇XYZ",
            "X,Y,Z 是正規化房間座標",
            "8 個係數對應 8 個角落感測器約束",
            "係數會依變數 v 與當次感測校正狀態而定",
        ],
        "為什麼剛好 8 點",
        [
            "房間有地面四角與天花板四角",
            "每個變數在同一時間有 8 個 residual",
            "三線性校正場也有 8 個自由度",
            "因此可由 8 個角點唯一決定此低階校正場",
            "但不是唯一決定任意複雜真實場",
        ],
    ),
    (
        "公式說明 18：角點 residual",
        "residual 定義",
        [
            "rᵛ_{abc}(t) = Oᵥ(p_{abc},t) - Nᵥ(p_{abc},t)",
            "a,b,c ∈ {0,1}",
            "p_{abc}：其中一個房間角點",
            "Oᵥ：該角點感測器觀測值",
            "Nᵥ：nominal model 在同一角點的預測值",
        ],
        "直覺意義",
        [
            "rᵛ 是主模型在感測點的誤差",
            "如果 rᵛ 為正，代表模型低估該角點",
            "如果 rᵛ 為負，代表模型高估該角點",
            "校正場 Cᵥ 的任務就是把這些角點誤差平滑帶入室內",
            "這一步讓模型與真實稀疏感測資料對齊",
        ],
    ),
    (
        "公式說明 19：三線性校正式",
        "校正公式",
        [
            "Cᵥ(X,Y,Z,t) = Σ_{a,b,c∈B} rᵛ_{abc}(t)",
            "× ℓ_a(X) ℓ_b(Y) ℓ_c(Z), B={0,1}",
            "ℓ₀(u)=1-u，ℓ₁(u)=u",
            "X/Y/Z 皆在 0 到 1 之間",
            "每個內部點都是 8 個角落 residual 的加權和",
        ],
        "重要性質",
        [
            "所有權重非負且總和為 1",
            "所以這是房間內部補間，不是無限制外插",
            "在任一角點上，對應權重為 1，其餘為 0",
            "因此校正後感測器位置會與觀測一致",
            "這是本研究 8 點推估最核心的數學基礎",
        ],
    ),
    (
        "公式說明 20：校正後估計值",
        "回到主公式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "在角點：Cᵥ 等於觀測 residual",
            "所以 F̂ᵥ(p_{abc},t) = Oᵥ(p_{abc},t)",
            "在非角點：Cᵥ 是 8 個 residual 的三線性補間",
            "Nᵥ 則保留設備與物理結構的空間趨勢",
        ],
        "適用範圍",
        [
            "8 顆感測器不能直接量到所有點",
            "三線性 residual correction 在角點與觀測一致",
            "其他點是 nominal model 加上低階 residual 補間的估計",
            "適用於主要空間趨勢估計",
            "不等同無條件還原任意真實室內場",
        ],
    ),
    (
        "公式說明 21：可完全表示的 residual 空間",
        "函數空間",
        [
            "𝒱 = span{1, X, Y, Z, XY, XZ, YZ, XYZ}",
            "這個空間的維度是 8",
            "三線性函數可由 8 個角點取值唯一決定",
            "如果真實 residual 屬於 𝒱",
            "則 8 個角點 residual 可完全重建整個 residual 場",
        ],
        "適用範圍",
        [
            "並非所有室內場都必然三線性",
            "三線性 residual 假設下可完全重建",
            "平滑但非三線性的 residual 可由誤差界描述接近程度",
            "突發局部熱源、光斑或遮蔽尖峰需額外資料補強",
            "此條件界定模型適用範圍",
        ],
    ),
    (
        "公式說明 22：平滑 residual 的誤差界",
        "這個上界在衡量什麼",
        [
            "|Rᵥ-Cᵥ| ≤ W²M_xx/8 + L²M_yy/8 + H²M_zz/8",
            "左邊：某個未量測點的 residual 可能補錯多少",
            "M_xx ≥ max|∂²Rᵥ/∂x²|：x 方向曲率上界",
            "M_yy、M_zz 同理對應 y/z 方向",
            "W/L/H 越大，補間跨度越長",
        ],
        "為什麼會這樣",
        [
            "xx 表示對 x 微分兩次，不是 x×x",
            "線性補間用端點連線近似中間曲線",
            "若 residual 是直線，二階導數為 0",
            "曲率越大，中間點越可能偏離直線",
            "三線性補間是 x/y/z 三方向線性補間",
            "平滑 residual 可控；尖峰或光斑需更多資料",
        ],
    ),
    (
        "公式說明 23：非連網裝置影響學習",
        "before/after delta",
        [
            "Δy_m = y_mᵃᶠᵗᵉʳ - y_mᵇᵉᶠᵒʳᵉ",
            "m ∈ {T,H,L}",
            "X_{i,k}：第 i 個感測點對第 k 個裝置 envelope",
            "Δy_m ≈ X β_m",
            "同一筆 record 綁定裝置狀態、baseline 與 outdoor",
        ],
        "least-squares 估計",
        [
            "β_m = argmin_{β} ||Δy_m - Xβ||²₂",
            "β_m：裝置對第 m 個因子的影響係數",
            "X：由 influence envelope 組成的設計矩陣",
            "before/after 觀測提供真實操作造成的差異訊號",
            "多事件重疊時係數代表混合效果",
        ],
    ),
    (
        "公式說明 24：Hybrid residual",
        "第二層修正",
        [
            "F_hybridᵛ(p,t) = Fᵥ(p,t) + Rᵥ(p,t; θᵥ)",
            "Fᵥ：前面可解釋的 base estimator",
            "Rᵥ：小型 neural network 預測的 residual",
            "預測推廣：ŷ_hybrid(t+h|I_t)=ŷ_phys(t+h|I_t)+ê(t+h|I_t)",
            "h 是 lead time；兩項對齊同一 target time",
        ],
        "定位",
        [
            "hybrid 不是取代物理模型",
            "physics 有 t+h：主模型須預測目標時刻 baseline",
            "I_t 只含 t 時可得資訊；禁止未來觀測 leakage",
            "本研究目前為 current-state spatial estimate，即 h=0",
            "LOO 結果代表標準情境 family 內 residual 可學習",
        ],
    ),
    (
        "公式說明 25：Hybrid 訓練目標",
        "residual label",
        [
            "Rᵥ*(p,t) = F_trueᵛ(p,t) - Fᵥ(p,t)",
            "F_trueᵛ：訓練或合成 truth 場",
            "Fᵥ：base estimator 輸出",
            "Rᵥ*：希望 neural network 學到的剩餘誤差",
            "訓練時不是直接預測整個場，而是預測 residual",
        ],
        "損失函數",
        [
            "ℒ(θᵥ) = (1/N)Σᵢ ||Rᵥ*(pᵢ,tᵢ)",
            "- Rᵥ(pᵢ,tᵢ;θᵥ)||² + λ||θᵥ||²",
            "第一項是 residual 預測誤差",
            "第二項是正則化，降低過擬合",
            "溫濕度響應較平滑可低通；照度因光源/遮蔽需保留快速跳變",
        ],
    ),
    (
        "公式說明 26：MAE、RMSE 與 Correlation",
        "誤差指標",
        [
            "MAE = (1/n) Σᵢ |ŷᵢ - yᵢ|",
            "RMSE = √[(1/n) Σᵢ (ŷᵢ - yᵢ)²]",
            "Corr = cov(ŷ,y)/(σ_{ŷ} σ_{y})",
            "ŷᵢ：模型估計值；yᵢ：truth 或觀測",
            "n：比較樣本數",
        ],
        "使用原因",
        [
            "MAE 代表平均偏差，最直觀",
            "RMSE 會放大尖峰或離群誤差",
            "Correlation 用於公開資料時序任務，檢查趨勢是否同向",
            "三者搭配可避免只看單一指標",
            "照度量級較大，圖表常用 log-scale 避免遮蔽溫濕度差異",
        ],
    ),
    (
        "公式說明 27：IDW baseline",
        "IDW 插值",
        [
            "IDW(p) = Σ_s w_s O_s / Σ_s w_s",
            "w_s = 1 / (dist(p,s) + ε)^q",
            "s：感測器索引",
            "O_s：感測器觀測值",
            "q：距離權重指數",
        ],
        "比較基準理由",
        [
            "IDW 是無設備物理先驗的幾何插值 baseline",
            "它只知道感測器位置與距離",
            "不知道冷氣出風、窗戶日照或燈具位置",
            "若本研究優於 IDW，表示設備與空間結構有提供額外資訊",
            "但 IDW 在無光源、平坦場時可能表現不差",
        ],
    ),
    (
        "公式說明 28：推薦排序與驗證",
        "推薦分數",
        [
            "q_m(S)= (1/K)Σ_k F_m(p_{k},t)",
            "P(q)=Σ_m w_m max(0,(|q_m-g_m|-δ_m)/δ_m)",
            "Score(a)=P(q_base)-P(q_a)",
            "m ∈ {T,H,L}",
            "Score 越高，預期改善越大",
        ],
        "驗證限制",
        [
            "S 是推薦評估的 sample scope 或目標區域",
            "g_m 是目標值，δ_m 是容許範圍",
            "目前推薦排序是 counterfactual simulation",
            "真正驗證需做 before/after intervention",
            "比較實際 comfort penalty reduction 與預測改善是否一致",
        ],
    ),
]

FORMULA_NUMERIC_EXAMPLES: dict[str, tuple[str, str]] = {
    "公式說明 1：三因子場與查詢點": (
        "p=(2,1,1.2), t=10 min → (T,H,L)=(27.4°C,60%,280 lux)",
        "若查詢點 p=(2,1,1.2)，時間 t=10 min，模型輸出 T=27.4°C、H=60%、L=280 lux，表示同一個座標與時間可同時取得三個環境量。",
    ),
    "公式說明 2：總估計式": (
        "N_T=27.0, C_T=-0.4 → F̂_T=26.6°C",
        "以溫度為例，若 nominal model 得到 N_T=27.0°C，角點 residual correction 給 C_T=-0.4°C，則校正後 F̂_T=27.0-0.4=26.6°C。",
    ),
    "公式說明 3：Indoor baseline": (
        "b₀=(29°C,67%,90 lux)",
        "一個房間的初始狀態可寫成 b₀=(T₀,H₀,L₀)=(29°C,67%,90 lux)，後續裝置影響都以這組基準值作為起點。",
    ),
    "公式說明 4：baseline 的取得方式": (
        "ΣT=232°C, |S|=8 → T₀=29°C",
        "若 8 顆角落感測器在 t_{ref} 的溫度總和為 232°C，則 T₀=232/8=29°C；濕度與照度可用相同平均方式取得。",
    ),
    "公式說明 5：高度正規化": (
        "z=1.2 m, Hᵣ=3 m → ζ=-0.1",
        "若房高 Hᵣ=3 m、查詢點高度 z=1.2 m，則 ζ=1.2/3-0.5=-0.1，代表位置略低於房間中高。",
    ),
    "公式說明 6：設備 activation": (
        "a=1, t=10, τ=10 → A=1-e^-1=0.632",
        "若穩態強度 a=1、時間常數 τ=10 min、啟動後 t=10 min，則 A(t)=1-exp(-10/10)=1-e^-1≈0.632。",
    ),
    "公式說明 7：influence envelope": (
        "A=0.8,R=0.5,D=0.9,V=1 → E=0.36",
        "若某點的時間強度 A=0.8、距離權重 R=0.5、方向權重 D=0.9、無遮蔽 V=1，則 E=0.8×0.5×0.9×1=0.36。",
    ),
    "公式說明 8：溫度場主式": (
        "29-1.2-0.5+0.2(-0.1)=27.28°C",
        "若 T₀=29°C、B_T=-1.2°C、S_T=-0.5°C、γ_T M(t)=0.2、ζ=-0.1，則 N_T=29-1.2-0.5+0.2×(-0.1)=27.28°C。",
    ),
    "公式說明 9：溫度的全室與局部項": (
        "B_T=-1.0+0.2+0.1=-0.7°C",
        "若冷氣全室項 -1.0°C、窗戶全室項 +0.2°C、燈具熱源 +0.1°C，則 B_T=-1.0+0.2+0.1=-0.7°C。",
    ),
    "公式說明 10：冷氣溫度項": (
        "s=-1,k=0.8,d=3,P=1,A=0.5 → B_ac,T=-1.2°C",
        "冷房模式 s_m=-1，若 k_ac,Tᵍ=0.8、d_T=3、P_ac=1、A_ac=0.5，則 B_ac,T=-1×0.8×3×1×0.5=-1.2°C。",
    ),
    "公式說明 11：窗戶與燈具溫度項": (
        "k=0.05,T_out-T₀=4 → B_win,T=+0.20°C",
        "若窗戶熱交換係數 k_win,Tᵍ=0.05，室外與室內基準溫差 T_out-T₀=4°C，且 P_win=A_win=1，則 B_win,T=0.05×4=+0.20°C。",
    ),
    "公式說明 12：濕度場主式": (
        "clip[0,100](67-5+1-0.2)=62.8%",
        "若 H₀=67%、B_H=-5%、S_H=+1%、γ_HMζ=0.2，則 N_H=clip[0,100](67-5+1-0.2)=62.8%。",
    ),
    "公式說明 13：濕度來源項": (
        "-0.4×10×0.5 + 0.02×8=-1.84%",
        "若冷氣除濕項為 -0.4×10×0.5=-2.0%，窗戶項因 H_out-H₀=8 而為 0.02×8=+0.16%，則全室濕度項合計 -1.84%。",
    ),
    "公式說明 14：照度場主式": (
        "max(0,90+250+120+40+30)=530 lux",
        "若 L₀=90 lux、窗戶直射 250 lux、燈具直射 120 lux、環境光 40 lux、一次反射 30 lux，則 N_L=max(0,90+250+120+40+30)=530 lux。",
    ),
    "公式說明 15：直射光與環境光": (
        "500×1×0.8×0.5×0.4×1=80 lux",
        "若燈具增益 G_light=500、P_light=1、A_light=0.8、方向權重 Φ=0.5、距離衰減 Q=0.4、可見性 V=1，則 L_lightᵈⁱʳ=80 lux。",
    ),
    "公式說明 16：一次漫反射": (
        "0.6×200×0.5×0.4×0.8×1=19.2 lux",
        "若某牆面 ρ=0.6、接收平均照度 200 lux、相對面積 0.5、距離衰減 0.4、方向 cosine 0.8、可見性 1，則一次反射貢獻約 19.2 lux。",
    ),
    "公式說明 17：8 參數校正多項式": (
        "X=Y=Z=0.5 → C=0.375",
        "若 c₀=0.2、c₁=0.4、c₂=-0.2、c₃=0.1、c₄=0.1，其餘係數為 0，且 X=Y=Z=0.5，則 C=0.2+0.4×0.5-0.2×0.5+0.1×0.5+0.1×0.25=0.375。",
    ),
    "公式說明 18：角點 residual": (
        "O_T=27.2,N_T=26.8 → r_T=+0.4°C",
        "若某角點感測器觀測溫度 O_T=27.2°C，而 nominal model 在同一點預測 N_T=26.8°C，則 residual r_T=27.2-26.8=+0.4°C。",
    ),
    "公式說明 19：三線性校正式": (
        "中心點權重各 1/8，Σr=1.6 → C=0.2",
        "在房間中心 X=Y=Z=0.5 時，8 個角點權重各為 1/8；若 8 個角點 residual 總和為 1.6，則 C=1.6/8=0.2。",
    ),
    "公式說明 20：校正後估計值": (
        "N_T=26.8,C_T=0.4 → F̂_T=27.2°C",
        "若 nominal model 在某點給 N_T=26.8°C，校正場給 C_T=+0.4°C，則校正後 F̂_T=26.8+0.4=27.2°C。",
    ),
    "公式說明 21：可完全表示的 residual 空間": (
        "R=0.2+0.3X+0.1Y，(0.5,0.5,0)→0.40",
        "若 residual 函數 R=0.2+0.3X+0.1Y，屬於三線性函數空間；在 X=0.5、Y=0.5、Z=0 時，R=0.2+0.15+0.05=0.40。",
    ),
    "公式說明 22：平滑 residual 的誤差界": (
        "W=6,L=4,H=3,M=(0.01,0.02,0.01) → bound≈0.096",
        "若 W=6、L=4、H=3，且 M_xx=0.01、M_yy=0.02、M_zz=0.01，則上界為 36×0.01/8 + 16×0.02/8 + 9×0.01/8 = 0.09625，約 0.096。",
    ),
    "公式說明 23：非連網裝置影響學習": (
        "X=[1,0.5], Δy=[-0.8,-0.4] → β=-0.8",
        "若兩個感測點的 envelope 為 X=[1,0.5]，觀測變化 Δy=[-0.8,-0.4]，單一係數 least squares 為 β=(XᵀΔy)/(XᵀX)=(-0.8-0.2)/(1+0.25)=-0.8。",
    ),
    "公式說明 24：Hybrid residual": (
        "F=27.0,R=-0.3 → F_hybrid=26.7°C",
        "若 base estimator 輸出 F=27.0°C，hybrid residual model 預測 R=-0.3°C，則 F_hybrid=27.0-0.3=26.7°C。若改寫為 h-step forecast，physics 與 residual 都必須對齊 t+h；h 是 lead time，不是 physics 參數，且 I_t 不得含 t+h 的實測真值。本研究現有空間估測等價於 h=0。",
    ),
    "公式說明 25：Hybrid 訓練目標": (
        "errors=(0.2,-0.1), λ||θ||²=0.01 → ℒ=0.035",
        "若兩筆 residual 預測誤差為 0.2 與 -0.1，平方平均為 (0.04+0.01)/2=0.025，再加上正則化 0.01，則 ℒ=0.035。",
    ),
    "公式說明 26：MAE、RMSE 與 Correlation": (
        "ŷ=[1,2,3], y=[1,2,4] → MAE=0.33, RMSE=0.58, Corr≈0.98",
        "若 ŷ=[1,2,3]、y=[1,2,4]，絕對誤差平均為 (0+0+1)/3=0.33，RMSE=sqrt(1/3)=0.58，相關係數約 0.98。",
    ),
    "公式說明 27：IDW baseline": (
        "O=(26,30), d=(1,3), q=2 → IDW≈26.4°C",
        "若兩個感測器讀值為 26°C 與 30°C，距離為 1 m 與 3 m，q=2，則權重為 1 與 1/9，IDW=(26+30/9)/(1+1/9)≈26.4°C。",
    ),
    "公式說明 28：推薦排序與驗證": (
        "q_base=30,g=26,δ=2 → P=1；q_a=27 → P=0；Score=1",
        "若只看溫度，目標 g=26°C、容許 δ=2°C；動作前 q_base=30°C 時 P=(|30-26|-2)/2=1，候選動作後 q_a=27°C 時 P=0，因此 Score=1-0=1。",
    ),
}


def add_formula_walkthrough(prs: Presentation, start_page: int, compact: bool = False) -> int:
    page = start_page
    for title, left_title, left_lines, right_title, right_lines in FORMULA_WALKTHROUGH:
        slide = new_slide(prs)
        add_title(slide, title, "公式、符號意義、假設與限制")
        card_height = 5.25 if compact else 5.35
        example = FORMULA_NUMERIC_EXAMPLES.get(title)
        right_lines_with_example = [*right_lines, f"數值例：{example[0]}"] if example else list(right_lines)
        add_card(slide, 0.65, 1.35, 6.05, card_height, left_title, left_lines, render_simple_subscripts=True)
        add_card(slide, 6.95, 1.35, 5.75, card_height, right_title, right_lines_with_example, render_simple_subscripts=True)
        add_footer(slide, page)
        page += 1
    return page


def build_presentation() -> Presentation:
    prs = init_presentation()
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    window_in_domain, window_out_of_domain = window_temperature_domain_counts(window_summary)
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    e8_summary = read_json(DATA / "e8_intervention_summary.json")
    rnn_summary = read_json(DATA / "public_benchmarks" / "rnn_sml2010_comparison.json")
    avg_mae = average_field_mae(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]
    bedroom_bootstrap = bedroom_aggregate["paired_day_block_bootstrap"]
    bedroom_lodo = bedroom_aggregate["leave_one_date_out_sensitivity"]
    bedroom_lodo_metrics = bedroom_lodo["metrics"]

    # Slide 1
    slide = new_slide(prs)
    add_title(
        slide,
        "單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型",
        "A Sparse-Sensing Spatial Digital Twin for Learning Environmental Impacts of Non-Networked Appliances in a Single Room",
    )
    add_bullets(
        slide,
        0.9,
        1.7,
        5.4,
        3.8,
        [
            "研究生：林昀佑",
            "指導教授：易昶霈教授、沈慧宇副教授",
            "系所：國立彰化師範大學資訊工程學系碩士班",
            "主題：以 8 顆角落感測器重建單房間溫度、濕度、照度場",
        ],
        level0_size=20,
    )
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 6.7, 1.5, 5.6, 4.8)
    add_footer(slide, 1)

    # Slide 2
    slide = new_slide(prs)
    add_title(slide, "研究問題與動機")
    add_card(
        slide,
        0.7,
        1.5,
        3.8,
        3.9,
        "問題背景",
        [
            "一般房間內仍存在大量非連網裝置",
            "冷氣、窗戶、照明會改變環境",
            "但系統通常無法直接讀到其狀態",
        ],
    )
    add_card(
        slide,
        4.75,
        1.5,
        3.8,
        3.9,
        "核心挑戰",
        [
            "只有少量感測器，無法直接知道全室分布",
            "裝置可能新增、移動，家具也會阻擋傳遞",
            "早期純插值與 local-only 模型都出現不合理結果",
        ],
    )
    add_card(
        slide,
        8.8,
        1.5,
        3.8,
        3.9,
        "研究目標",
        [
            "重建三因子空間場",
            "學習非連網裝置對環境的影響",
            "以 Web 與工具介面提供查詢與決策能力",
        ],
    )
    add_footer(slide, 2)

    # Slide 3
    slide = new_slide(prs)
    add_title(slide, "論文整體邏輯：問題、方法、證據與結論邊界")
    add_picture(slide, ARCHITECTURE / "研究整體邏輯架構.svg", 0.55, 1.25, 12.25, 5.75)
    add_footer(slide, 3)

    # Slide 4
    slide = new_slide(prs)
    add_title(slide, "房間拓樸、感測器與目標區域")
    add_picture(slide, ARCHITECTURE / "房間感測器與目標區域配置.svg", 0.8, 1.4, 6.0, 5.3)
    add_card(
        slide,
        7.1,
        1.55,
        5.0,
        4.6,
        "固定設定",
        [
            "房間尺寸：6 m × 4 m × 3 m",
            "感測器：地面四角 + 天花板四角，共 8 顆",
            "區域：window_zone / center_zone / door_side_zone",
            "設備：ac_main / window_main / light_main",
            "採樣網格：16 × 12 × 6",
        ],
    )
    add_footer(slide, 4)

    # Slide 5
    slide = new_slide(prs)
    add_title(slide, "數學模型：變數專屬場模型")
    add_card(
        slide,
        0.7,
        1.45,
        5.9,
        4.8,
        "核心表示式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "Nᵥ：變數專屬 nominal model",
            "Cᵥ：由角落感測器殘差擬合的 trilinear 校正場",
            "溫度：熱交換、熱源、垂直分層",
            "濕度：除濕、水氣交換、物理範圍限制",
            "照度：直射、環境光、一次漫反射",
        ],
    )
    add_card(
        slide,
        6.85,
        1.45,
        5.8,
        4.8,
        "模型特點",
        [
            "冷氣、窗戶、照明都是模組化裝置",
            "家具會造成 obstacle-aware attenuation",
            "動態響應採一階收斂近似",
            "三個變數共用座標、設備與校正框架",
            "但 nominal model 不共用同一套物理公式",
            "推薦必須先有 sample / cluster 與三因子目標",
        ],
    )
    add_footer(slide, 5)

    # Slide 6
    slide = new_slide(prs)
    add_title(slide, "模型學習、推論與推薦資料流")
    add_picture(slide, ARCHITECTURE / "模型學習推論與推薦資料流.svg", 0.45, 1.12, 12.45, 5.72)
    add_footer(slide, 6)

    # Slide 7
    slide = new_slide(prs)
    add_title(slide, "系統實作與介面")
    add_card(
        slide,
        0.7,
        1.4,
        3.8,
        4.9,
        "MCP：工具化介面",
        [
            "MCP 不做預測；只提供 AI 可呼叫 runtime tools",
            "initialize：設定 scenario、baseline、外部邊界、設備/家具、時間與 estimator",
            "sample point：查指定座標在特定時間/穩定態的三因子",
            "learn impacts：建立 before/after 觀測紀錄再學係數",
            "window direct：直接輸入外部溫濕度、日照與開窗比例",
            "rank actions：指定座標 sample + T/H/L 目標後才排序註冊設備操作",
        ],
    )
    add_card(
        slide,
        4.75,
        1.4,
        3.8,
        4.9,
        "Web Demo",
        [
            "可旋轉 3D 預覽",
            "時間軸、播放、point sample",
            "裝置與家具都可模組化設定",
            "支援 hybrid estimator toggle",
        ],
    )
    add_card(
        slide,
        8.8,
        1.4,
        3.8,
        4.9,
        "輸入模式",
        [
            "標準情境 8 組",
            f"窗戶矩陣 {window_summary.get('count', 0)} 組：{window_in_domain} 範圍內／{window_out_of_domain} 壓力測試",
            "窗戶 direct input",
            "自訂家具與自訂裝置",
        ],
    )
    add_footer(slide, 7)

    add_learning_record_slide(prs, 8)

    # Slide 8
    slide = new_slide(prs)
    add_title(slide, "驗證流程與比較原則")
    add_picture(slide, ARCHITECTURE / "驗證與實驗流程圖.svg", 0.8, 1.35, 5.8, 5.1)
    add_bullets(
        slide,
        6.95,
        1.5,
        5.25,
        5.2,
        [
            "E1-E3：synthetic full-field、IDW baseline、ablation",
            "E4：非連網裝置影響學習與推薦排序",
            f"E5：window matrix {window_summary.get('count', 0)} 組（{window_in_domain} 範圍內／{window_out_of_domain} 壓力測試）",
            "E6：hybrid no-Fourier 對照與 LOO cross-validation",
            f"E7：bedroom_01 {bedroom_summary['snapshot_count']} 筆，pillow hold-out",
            f"E8 execution kit：schema / template / analyzer；{e8_summary['trial_counts']['completed']} trials、{e8_summary['evidence_status']}",
            "E9 public benchmark；demo 不是量化實驗",
        ],
        level0_size=17,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = new_slide(prs)
    add_title(slide, "主要結果：場重建與 baseline 比較")
    add_card(
        slide,
        0.7,
        1.45,
        3.7,
        2.0,
        "平均 Field MAE",
        [
            f"Temperature: {avg_mae['temperature']}",
            f"Humidity: {avg_mae['humidity']}",
            f"Illuminance: {avg_mae['illuminance']}",
        ],
    )
    add_card(
        slide,
        0.7,
        3.75,
        3.7,
        2.3,
        "標準情境",
        [
            f"驗證情境數：{len(validation_summary.get('scenarios', []))}",
            "比較項目：field / sensors / zones / IDW / recommendations",
        ],
    )
    add_picture(slide, FIGURES / "all_active_temperature_3d.svg", 4.7, 1.45, 3.8, 4.7)
    add_picture(slide, FIGURES / "window_only_illuminance_3d.svg", 8.7, 1.45, 3.8, 4.7)
    add_footer(slide, 9)

    # Slide 10
    slide = new_slide(prs)
    add_title(slide, "Hybrid Residual Neural Network 結果")
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]
    add_card(
        slide,
        0.7,
        1.45,
        5.0,
        4.9,
        "Held-out + LOO",
        [
            f"Default samples: {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}",
            f"Default hybrid MAE: {metric_triplet(default_hybrid['hybrid_test_field_mae'])}",
            f"No-Fourier hybrid MAE: {metric_triplet(no_fourier['hybrid_test_field_mae'])}",
            f"LOO avg hybrid MAE: {metric_triplet(loo['average_hybrid_field_mae'])}",
            f"LOO reduction: T {loo['average_field_mae_reduction_percent']['temperature']:.2f}%, H {loo['average_field_mae_reduction_percent']['humidity']:.2f}%, L {loo['average_field_mae_reduction_percent']['illuminance']:.2f}%",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 6.0, 1.35, 6.2, 3.1)
    add_bullets(
        slide,
        6.1,
        4.65,
        5.9,
        1.5,
        [
            "hybrid residual 是第二層修正器，不取代主模型",
            "LOO 結果證明標準情境 family 內殘差可學習，不代表任意房間泛化",
            f"E7 以 {bedroom_bootstrap['replicates']:,} 次 date-block bootstrap 檢查 pillow 改善，三因子 95% CI 下界皆 > 0",
            f"E7 逐日剔除最小降幅：T {bedroom_lodo_metrics['temperature']['minimum_absolute_mae_reduction']:.4f} / H {bedroom_lodo_metrics['humidity']['minimum_absolute_mae_reduction']:.4f} / L {bedroom_lodo_metrics['illuminance']['minimum_absolute_mae_reduction']:.4f}",
        ],
        level0_size=13,
    )
    add_footer(slide, 10)

    # Slide 11
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：SML2010 / CU-BEMS")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "sml2010_task_breakdown.svg", 0.6, 1.35, 6.0, 3.35)
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "cu_bems_task_breakdown.svg", 6.75, 1.35, 6.0, 3.35)
    add_bullets(
        slide,
        0.9,
        5.05,
        11.6,
        1.6,
        [
            "S3 是主要優勢：事件/邊界 delta 需要變化方向，structured prior 比 persistence 與 linear regression 更有用",
            "S1 與 C2 是主要劣勢：短視窗照度高度自相關，且公開資料缺實際幾何、遮蔽與多燈具資訊",
            "Oh2024-inspired residual：15min 兩點溫度最佳；60min 本研究 readout 最佳；24h persistence 最佳",
            "次日 primary：validation 選出的 daily trend 反而惡化 7.34% / 8.36%；bootstrap interval 均跨 0",
            "探索性 adaptive median 亦惡化 8.83% / 9.73%；約 1% 的未選中 bias correction 不足以主張優勢",
            f"RNN 同資料比較：{rnn_summary['summary']['evaluated_cases']}/12 通過 parity；lowest MAE = sequence LR 7、persistence 5、RNN 0",
            "四種方法共用四筆 history、split、targets 與 test rows；負向結果保留",
        ],
        level0_size=13,
    )
    add_footer(slide, 11)

    # Slide 12
    slide = new_slide(prs)
    add_title(slide, "研究貢獻與資料策略")
    add_bullets(
        slide,
        0.9,
        1.55,
        11.4,
        5.3,
        [
            "提出以單房間、8 顆角落感測器為前提的三因子空間數位孿生原型",
            "以變數專屬 nominal model、power calibration 與 trilinear correction 建立可解釋估測流程",
            "以 least-squares 學習非連網裝置影響，並用 hybrid residual 做第二層修正",
            "明確拆分 synthetic full-field、real sparse calibration、public task-aligned 與 intervention validation",
            "完整 3D 場比較以 canonical synthetic benchmark 為主",
            f"真實臥室快照校正後 pillow MAE: {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            "公開資料集則採 task-aligned benchmark：CU-BEMS / SML2010 / ASHRAE 各比相容子任務",
            "室內應用溫度範圍固定 20–30 °C；人體舒適採目標帶與容許範圍",
        ],
        level0_size=16,
    )
    add_footer(slide, 12)

    # Slide 13
    slide = new_slide(prs)
    add_title(slide, "結論與未來工作")
    add_card(
        slide,
        0.8,
        1.6,
        5.7,
        4.8,
        "結論",
        [
            "有限角落感測器下仍可用分層式模型重建單房間三因子分布",
            "非連網裝置可透過環境變化進行影響學習與校正",
            "bedroom_01 7 天快照顯示校正後可改善未參與 fitting 的 pillow 點",
            "各資料來源支援的驗證範圍已拆開說明",
            "模型已能輸出區域估計、反事實推薦排序與 AI 可查詢工具",
        ],
    )
    add_card(
        slide,
        6.8,
        1.6,
        5.7,
        4.8,
        "未來工作",
        [
            "擴大 ESP32 長期真實資料",
            "擴充 CO2 / PM2.5 等因子",
            "改進 multi-zone / partition 模型",
            "補足 dense real-room ground truth",
            "執行推薦動作 before/after 介入驗證",
            "20–30 °C 動態植物生長環境：先補 PPFD/CO2/基質與生物 endpoint",
            "定義 state/observation/noise 後比較 moving average、KF 與 EKF",
        ],
    )
    add_footer(slide, 13)
    add_formula_walkthrough(prs, 14, compact=True)
    renumber_footers(prs)
    return prs


def build_presentation_30min() -> Presentation:
    prs = init_presentation()
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    window_in_domain, window_out_of_domain = window_temperature_domain_counts(window_summary)
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    e8_summary = read_json(DATA / "e8_intervention_summary.json")
    rnn_summary = read_json(DATA / "public_benchmarks" / "rnn_sml2010_comparison.json")
    avg_mae = average_field_mae(validation_summary)
    scenarios = scenario_map(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]
    bedroom_bootstrap = bedroom_aggregate["paired_day_block_bootstrap"]
    bedroom_bootstrap_metrics = bedroom_bootstrap["metrics"]
    bedroom_lodo = bedroom_aggregate["leave_one_date_out_sensitivity"]
    bedroom_lodo_metrics = bedroom_lodo["metrics"]

    # 1 cover
    slide = new_slide(prs)
    add_title(
        slide,
        "單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型",
        "碩士論文簡報",
    )
    add_bullets(
        slide,
        0.9,
        1.7,
        5.4,
        4.0,
        [
            "研究生：林昀佑",
            "指導教授：易昶霈教授、沈慧宇副教授",
            "國立彰化師範大學資訊工程學系碩士班",
            "主題：單房間三因子數位孿生、非連網裝置影響學習、工具化服務介面",
        ],
        level0_size=20,
    )
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 6.7, 1.5, 5.6, 4.8)
    add_footer(slide, 1)

    # 2 roadmap
    slide = new_slide(prs)
    add_title(slide, "報告流程")
    add_bullets(
        slide,
        1.0,
        1.55,
        11.0,
        5.4,
        [
            "1. 問題背景與研究動機",
            "2. 文獻定位與研究缺口",
            "3. 系統架構與數學模型",
            "4. 感測器校正、影響學習與 hybrid residual",
            "5. 實作系統、驗證設計與實驗結果",
            "6. 結論、限制、未來工作與公式整理",
        ],
        level0_size=21,
    )
    add_footer(slide, 2)

    # research logic overview
    slide = new_slide(prs)
    add_title(slide, "論文整體邏輯：問題、方法、證據與結論邊界")
    add_picture(slide, ARCHITECTURE / "研究整體邏輯架構.svg", 0.55, 1.25, 12.25, 5.75)
    add_footer(slide, 3)

    # 3 background
    slide = new_slide(prs)
    add_title(slide, "研究背景與問題")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        4.8,
        "房間場景",
        [
            "真實房間通常只有少量感測器",
            "但使用者關心的是整個空間的舒適度",
            "不是單一點位數值",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        4.8,
        "非連網裝置",
        [
            "冷氣、窗戶、照明常無 API",
            "狀態不可直接讀取",
            "卻持續改變溫度、濕度、照度",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.8,
        "核心需求",
        [
            "重建全室三因子分布",
            "學習裝置對環境的影響",
            "在 sample/cluster + 三因子目標下支援推薦",
        ],
    )
    add_footer(slide, 4)

    # 4 questions and contributions
    slide = new_slide(prs)
    add_title(slide, "研究問題與貢獻")
    add_card(
        slide,
        0.75,
        1.45,
        5.6,
        4.9,
        "研究問題",
        [
            "RQ1：8 顆角落感測器能否重建單房間三因子空間場？",
            "RQ2：能否從環境資料學習非連網裝置影響？",
            "RQ3：sample/cluster 與三因子目標下能否排序控制動作？",
            "RQ4：能否將模型封裝成 MCP 可查詢工具？",
        ],
    )
    add_card(
        slide,
        6.7,
        1.45,
        5.9,
        4.9,
        "主要貢獻",
        [
            "single-room three-factor spatial digital twin",
            "power calibration + trilinear residual correction",
            "non-networked appliance impact learning",
            "hybrid residual + Fourier denoising",
            "task-aligned public benchmark strategy",
            "MCP / Gemma / Web 可互動原型",
        ],
    )
    add_footer(slide, 4)

    # 5 literature gap
    slide = new_slide(prs)
    add_title(slide, "文獻定位、研究缺口與比較原則")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        5.0,
        "已有研究",
        [
            "房間尺度 IEQ 實驗",
            "有限感測器場重建",
            "hybrid thermal model",
            "建築 digital twin 平台",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        5.0,
        "常見不足",
        [
            "多只看熱環境或單雙因子",
            "少處理非連網裝置學習",
            "少將模型做成 AI 可查詢工具",
            "多半不是單房間低成本原型",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        5.0,
        "本研究定位",
        [
            "single-room",
            "limited corner sensors",
            "temperature + humidity + illuminance",
            "control-oriented + MCP-accessible",
            "public datasets only for aligned subtasks",
        ],
    )
    add_footer(slide, 5)

    # 6 architecture
    slide = new_slide(prs)
    add_title(slide, "整體系統架構")
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 0.8, 1.35, 6.2, 5.3)
    add_bullets(
        slide,
        7.15,
        1.5,
        5.0,
        5.2,
        [
            "圖 3-1 以 top-down tree 呈現系統責任邊界",
            "情境與觀測層提供 room schema、8 點感測、外部邊界與時間",
            "估測與學習層負責三因子場模型、校正、影響學習與 hybrid residual",
            "服務與決策層暴露 scripts、Web、MCP/Gemma 與 point/zone/action 輸出",
        ],
        level0_size=17,
    )
    add_footer(slide, 6)

    # 7 execution flow
    slide = new_slide(prs)
    add_title(slide, "主要執行資料流")
    add_picture(slide, ARCHITECTURE / "主要執行資料流.svg", 0.8, 1.35, 5.9, 5.25)
    add_bullets(
        slide,
        6.95,
        1.45,
        5.25,
        5.3,
        [
            "scenario 或 direct input 先進入 service",
            "套用 indoor baseline、裝置、家具與時間設定",
            "再做場估測、感測器校正與 dashboard 輸出",
            "MCP 和 Web 都走同一條執行路徑",
        ],
        level0_size=17,
    )
    add_footer(slide, 7)

    # 8 room topology
    slide = new_slide(prs)
    add_title(slide, "房間拓樸、感測器與目標區域")
    add_picture(slide, ARCHITECTURE / "房間感測器與目標區域配置.svg", 0.8, 1.35, 6.0, 5.3)
    add_card(
        slide,
        7.05,
        1.45,
        5.15,
        4.9,
        "固定研究設定",
        [
            "房間：6 × 4 × 3 m",
            "感測器：floor/ceiling 四角，共 8 顆",
            "區域：window / center / door-side",
            "核心裝置：ac_main / window_main / light_main",
            "解析度：16 × 12 × 6",
        ],
    )
    add_footer(slide, 8)

    # 9 devices and furniture
    slide = new_slide(prs)
    add_title(slide, "模組化裝置與家具阻擋")
    add_picture(slide, ARCHITECTURE / "可模組化裝置與家具架構.svg", 0.8, 1.35, 5.9, 5.1)
    add_bullets(
        slide,
        6.95,
        1.45,
        5.25,
        5.25,
        [
            "冷氣、窗戶、燈都可視為模組化裝置",
            "家具是可開關、可移動的阻擋物件",
            "阻擋效果會依幾何位置自適應調整",
            "Web 端可新增 custom devices 與 custom furniture",
        ],
        level0_size=17,
    )
    add_footer(slide, 9)

    # 10 math model
    slide = new_slide(prs)
    add_title(slide, "數學模型：變數專屬 nominal model")
    add_card(
        slide,
        0.75,
        1.4,
        6.0,
        5.0,
        "核心形式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "Nᵥ：溫度、濕度、照度各自的 nominal model",
            "Cᵥ：8 顆角落感測器建立的 residual correction",
            "溫度處理熱交換與垂直分層",
            "濕度處理除濕與外氣水氣交換",
            "照度處理光源幾何、遮蔽與反射",
        ],
    )
    add_card(
        slide,
        7.0,
        1.4,
        5.2,
        5.0,
        "設計理由",
        [
            "避免只有冷氣附近變冷、全室仍近乎不變的不合理情況",
            "早期純插值與 local-only 版本都失敗過",
            "避免把同一組 bulk/local 公式硬套到三種物理量",
            "比 CFD 輕量",
            "比純插值更可解釋",
            "適合控制導向與即時服務化",
        ],
    )
    add_footer(slide, 10)

    # method rationale
    slide = new_slide(prs)
    add_title(slide, "方法選擇：為什麼不是純插值、純物理或純黑盒")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.95,
        "純插值不足",
        [
            "IDW 只看距離與感測點數值",
            "不知道冷氣出風口、窗戶方向、燈具位置",
            "照度與局部熱區容易被抹平",
            "適合作 baseline，不適合作主模型",
        ],
        body_size=13,
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.95,
        "完整物理太重",
        [
            "CFD / ray tracing 需要更多邊界條件與計算成本",
            "一般房間很難取得精確材質、風場與設備曲線",
            "本研究目標是控制導向與即時服務化",
            "因此採 reduced-order nominal model",
        ],
        body_size=13,
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.95,
        "純黑盒風險",
        [
            "小資料下容易過擬合",
            "難解釋設備、邊界與空間位置的作用",
            "hybrid residual 只學剩餘誤差",
            "保留主模型的可解釋結構",
        ],
        body_size=13,
    )
    add_footer(slide, 11)

    # 11 calibration and learning
    slide = new_slide(prs)
    add_title(slide, "模型學習、推論與推薦資料流")
    add_picture(slide, ARCHITECTURE / "模型學習推論與推薦資料流.svg", 0.45, 1.1, 12.45, 5.75)
    add_footer(slide, 11)

    # 12 implementation interfaces
    slide = new_slide(prs)
    add_title(slide, "系統實作與介面")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.9,
        "MCP：工具化介面",
        [
            "MCP 不做預測；核心模型負責場估計、校正與排序",
            "initialize：註冊 scenario、baseline、外部邊界、設備/家具、時間與 estimator",
            "AC state：模式、目標溫度、風量、水平/垂直角度與擺動",
            "sample point：補足非感測點、可指定 elapsed/steady state",
            "learn impacts：start/finish before-after record",
            "window direct：輸入外部資料，不走 48 組 preset",
            "rank actions：以指定座標 sample 與 T/H/L 目標評估註冊設備操作",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.9,
        "Gemma / Ollama",
        [
            "Gemma 透過 bridge 做 tool calling",
            "MCP 支援來自主機與 runtime",
            "不是宣稱模型原生支援 MCP",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "Web Demo",
        [
            "可旋轉 3D 預覽",
            "時間軸與播放",
            "裝置/家具模組化調整",
            "hybrid estimator toggle",
        ],
    )
    add_footer(slide, 12)

    add_learning_record_slide(prs, 13)

    # 13 validation design
    slide = new_slide(prs)
    add_title(slide, "驗證設計")
    add_picture(slide, ARCHITECTURE / "驗證與實驗流程圖.svg", 0.75, 1.35, 6.0, 5.2)
    add_bullets(
        slide,
        7.0,
        1.45,
        5.2,
        5.3,
        [
            "E1-E3：truth-adjusted simulation、IDW、synthetic ablation",
            f"E4-E6：裝置影響學習、{window_summary.get('count', 0)} 組 window matrix、hybrid no-Fourier/LOO",
            f"E7：bedroom_01 {bedroom_summary['snapshot_count']} 筆快照與 pillow 位置比較",
            "E8：推薦動作 before/after intervention protocol",
            "E9：public datasets 僅作 task-aligned benchmark",
            "Web demo 與 3D 展示是呈現層，不列為量化實驗",
        ],
        level0_size=16,
    )
    add_footer(slide, 13)

    # evidence boundary
    slide = new_slide(prs)
    add_title(slide, "證據鏈與驗證範圍")
    add_card(
        slide,
        0.65,
        1.35,
        3.0,
        5.25,
        "Synthetic full-field",
        [
            "支援完整 3D 場誤差比較",
            "可比較 IDW、base、ablation、hybrid",
            "限制：truth 仍來自受控生成與調整，不等同長期真實場",
        ],
        body_size=12,
    )
    add_card(
        slide,
        3.85,
        1.35,
        3.0,
        5.25,
        "Real-bedroom snapshot",
        [
            "支援稀疏感測校正的真實點位檢查",
            "pillow hold-out 不參與 fitting",
            "限制：不是 dense real-room ground truth",
        ],
        body_size=12,
    )
    add_card(
        slide,
        7.05,
        1.35,
        2.85,
        5.25,
        "Public datasets",
        [
            "支援相容子任務比較",
            "SML2010 / CU-BEMS 只做 task-aligned benchmark",
            "限制：缺單房間幾何與 8 點拓樸",
        ],
        body_size=12,
    )
    add_card(
        slide,
        10.1,
        1.35,
        2.7,
        5.25,
        "Recommendation",
        [
            "目前可做反事實排序",
            "E8 定義介入驗證 protocol",
            "限制：尚未完成真實 before/after 因果驗證",
        ],
        body_size=12,
    )
    add_footer(slide, 14)

    # 14 scenarios and time/window settings
    slide = new_slide(prs)
    add_title(slide, "情境設計與輸入模式")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        4.9,
        "標準情境",
        [
            "idle",
            "ac_only / window_only / light_only",
            "ac_window / window_light / ac_light",
            "all_active",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        4.9,
        "窗戶模式",
        [
            "四季 × 天氣 × 時段",
            f"{window_in_domain} 組範圍內／{window_out_of_domain} 組範圍外壓力測試",
            "也支援 direct outdoor input",
            "可分析窗邊區與中心區差異",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "時間軸",
        [
            "所有 scenario 都有 elapsed time",
            "近似一階動態收斂",
            "Web 端可播放到 quasi-steady state",
        ],
    )
    add_footer(slide, 14)

    # 15 quantitative results
    slide = new_slide(prs)
    add_title(slide, "主要量化結果")
    add_card(
        slide,
        0.7,
        1.4,
        3.6,
        2.0,
        "平均 Field MAE",
        [
            f"Temperature: {avg_mae['temperature']}",
            f"Humidity: {avg_mae['humidity']}",
            f"Illuminance: {avg_mae['illuminance']}",
        ],
    )
    add_card(
        slide,
        0.7,
        3.7,
        3.6,
        2.4,
        "圖表資料",
        [
            "8 scenarios, full 3D grid",
            "Y axis: log-scale Field MAE",
            "Bars: IDW / Base / LOO Hybrid",
            "IDW uses same 8 corner readings",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 4.55, 1.25, 7.7, 3.25)
    add_card(
        slide,
        4.6,
        4.65,
        7.7,
        1.7,
        "真實臥室校正檢查",
        [
            f"Raw pillow MAE: {metric_triplet(bedroom_aggregate['raw_pillow_mae'])}",
            f"Corrected pillow MAE: {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            f"{bedroom_bootstrap['replicates']:,} 次 date-block bootstrap：三因子 MAE 降幅 95% CI 下界皆 > 0",
            f"逐日剔除最小降幅 T/H/L：{bedroom_lodo_metrics['temperature']['minimum_absolute_mae_reduction']:.4f} / {bedroom_lodo_metrics['humidity']['minimum_absolute_mae_reduction']:.4f} / {bedroom_lodo_metrics['illuminance']['minimum_absolute_mae_reduction']:.4f}",
        ],
        body_size=11,
    )
    add_footer(slide, 15)

    # real-room and recommendation status
    slide = new_slide(prs)
    add_title(slide, "真實臥室快照與推薦驗證狀態")
    add_card(
        slide,
        0.7,
        1.45,
        5.75,
        4.95,
        "E7：real-bedroom sparse calibration",
        [
            f"{bedroom_summary['snapshot_count']} snapshots；pillow 不參與 8 角點 fitting",
            f"Raw → corrected MAE: {metric_triplet(bedroom_aggregate['raw_pillow_mae'])} → {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            f"Paired bootstrap：{bedroom_bootstrap['replicates']:,} 次，以 7 個 date blocks 重抽樣",
            "95% CI (T/H/L reduction):",
            f"[{bedroom_bootstrap_metrics['temperature']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['temperature']['ci95_absolute_mae_reduction']['upper']:.4f}] / [{bedroom_bootstrap_metrics['humidity']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['humidity']['ci95_absolute_mae_reduction']['upper']:.4f}] / [{bedroom_bootstrap_metrics['illuminance']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['illuminance']['ci95_absolute_mae_reduction']['upper']:.4f}]",
            f"改善快照：T {bedroom_bootstrap_metrics['temperature']['snapshots_improved']}/28；H/L 28/28；仍非 dense truth",
            f"7-fold 逐日剔除最小降幅：T {bedroom_lodo_metrics['temperature']['minimum_absolute_mae_reduction']:.4f}；H {bedroom_lodo_metrics['humidity']['minimum_absolute_mae_reduction']:.4f}；L {bedroom_lodo_metrics['illuminance']['minimum_absolute_mae_reduction']:.4f}",
        ],
        body_size=10,
    )
    add_card(
        slide,
        6.8,
        1.45,
        5.75,
        4.95,
        "E8：recommendation validation protocol",
        [
            "目前 rank actions 是模型反事實重跑後的 penalty reduction 排序",
            "versioned schema、空白 trial template 與 deterministic analyzer 已完成",
            f"完成真實介入 trials：{e8_summary['trial_counts']['completed']}；status：{e8_summary['evidence_status']}",
            "所有 efficacy estimates 維持 null；synthetic tests 只驗證公式",
            "正式驗證仍須 before/after intervention 與 matched controls",
        ],
        body_size=12,
    )
    add_footer(slide, 16)

    # 16 qualitative visual results
    slide = new_slide(prs)
    add_title(slide, "3D 視覺化結果")
    add_picture(slide, FIGURES / "all_active_temperature_3d.svg", 0.65, 1.5, 4.0, 4.8)
    add_picture(slide, FIGURES / "window_only_illuminance_3d.svg", 4.68, 1.5, 4.0, 4.8)
    add_picture(slide, FIGURES / "ac_only_temperature_3d.svg", 8.7, 1.5, 4.0, 4.8)
    add_footer(slide, 16)

    # 17 hybrid result
    slide = new_slide(prs)
    add_title(slide, "Hybrid Residual Neural Network 結果")
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]
    add_card(
        slide,
        0.75,
        1.45,
        5.2,
        4.9,
        "Robustness checks",
        [
            f"Default samples: {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}",
            f"Default MAE: {metric_triplet(default_hybrid['hybrid_test_field_mae'])}",
            f"No-Fourier MAE: {metric_triplet(no_fourier['hybrid_test_field_mae'])}",
            f"LOO avg MAE: {metric_triplet(loo['average_hybrid_field_mae'])}",
            f"LOO reduction: T {loo['average_field_mae_reduction_percent']['temperature']:.2f}%, H {loo['average_field_mae_reduction_percent']['humidity']:.2f}%, L {loo['average_field_mae_reduction_percent']['illuminance']:.2f}%",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 6.15, 1.35, 6.1, 3.0)
    add_bullets(
        slide,
        6.25,
        4.55,
        5.8,
        1.7,
        [
            "no-Fourier 對照顯示照度改善不是頻域處理造成",
            "LOO 降低單一 held-out split 過度樂觀的風險",
            "LOO 結果仍限標準情境 family，不等同任意房間泛化",
            "真實臥室快照已驗證 calibration，推薦有效性仍需介入實驗",
        ],
        level0_size=15,
    )
    add_footer(slide, 17)

    # 18 SML2010 public task breakdown
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：SML2010")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "sml2010_task_breakdown.svg", 0.75, 1.3, 7.1, 4.0)
    add_card(
        slide,
        8.15,
        1.4,
        4.35,
        4.85,
        "S1 / S2 / S3 判讀",
        [
            "SML2010 24 任務：12 lowest MAE、15 勝 LR、14 勝 persistence",
            "原 E9：S1 照度弱；S2 混合；S3 event delta 最強",
            "Oh-inspired：15min 兩點 T 最佳；本研究 readout 在 60min 最佳",
            "24h 兩點 T 皆由 persistence 最佳；transfer 亦劣於 raw physics",
            "次日 primary 選中 trend 但 test 惡化 7.34% / 8.36%，CI 均跨 0",
            "Post-primary adaptive 亦惡化；未建立 next-day advantage",
            f"RNN 公平比較：{rnn_summary['summary']['evaluated_cases']}/12 parity 通過；lowest MAE 為 sequence LR 7、persistence 5、RNN 0",
            "資料 confidential；只稱 method transfer，不稱原文重現或 full 3D 驗證",
        ],
        body_size=11,
    )
    add_footer(slide, 18)

    # 19 CU-BEMS public task breakdown
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：CU-BEMS")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "cu_bems_task_breakdown.svg", 0.75, 1.3, 7.1, 4.0)
    add_card(
        slide,
        8.15,
        1.4,
        4.35,
        4.85,
        "C1 / C2 / C3 判讀",
        [
            "C1：AC 溫濕度可補強 LR，但不勝 persistence",
            "C2：商辦照度與單房間假設差距大，是劣勢",
            "C3：compound event 可穩定勝過 LR",
            "CU-BEMS 12 任務：9 勝 LR、0 勝 persistence",
            "結果是 zone-level 外部壓力測試，不是 8 點房間拓樸驗證",
        ],
    )
    add_footer(slide, 19)

    # 20 conclusion and future
    slide = new_slide(prs)
    add_title(slide, "結論、限制與未來工作")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.9,
        "結論",
        [
            "單房間三因子數位孿生原型已可運作",
            "可估場、可校正、可學習、可推薦",
            "bedroom_01 快照顯示校正後 pillow 點 MAE 明顯下降",
            "資料比較需依任務層級切分",
            "hybrid residual 泛化仍需更多房間與 dense ground truth",
            "MCP 與 Web 展示已完成",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.9,
        "限制",
        [
            "已有小型真實臥室快照，但仍缺長期 dense field",
            "LOO hybrid 目前只支持標準情境 family 內殘差可學習",
            "不是 CFD 等級模型",
            "公開資料集缺乏 full-field ground truth",
            "室內溫度適用範圍限 20–30 °C；人體舒適以 tolerance 判定",
            "固定 RNN 同資料比較 lowest MAE 為 0/12",
            "推薦動作尚未完成真實介入式因果驗證",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "未來工作",
        [
            "擴大 ESP32 長期真實資料",
            "加入 CO2 / PM2.5",
            "發展 multi-zone / partition model",
            "補足 PPFD/CO2/基質與生物 endpoint 後再評估動態植物生長情境",
            "定義狀態與雜訊模型後比較 moving average、KF、EKF",
            "執行推薦動作 before/after 介入驗證",
            "朝閉環控制延伸",
        ],
    )
    add_footer(slide, 20)

    # formula section guide
    slide = new_slide(prs)
    add_title(slide, "公式與指標整理")
    add_card(
        slide,
        0.7,
        1.4,
        3.85,
        4.95,
        "場模型",
        [
            "三因子場與查詢點",
            "總估計式",
            "indoor baseline",
            "高度正規化",
            "設備 activation 與 envelope",
        ],
        body_size=13,
    )
    add_card(
        slide,
        4.75,
        1.4,
        3.85,
        4.95,
        "三因子公式",
        [
            "溫度：全室熱響應 + 局部熱場",
            "濕度：除濕 + 外氣水氣交換",
            "照度：直射、環境光、一次漫反射",
            "重點是三種物理量不共用同一套 nominal model",
        ],
        body_size=13,
    )
    add_card(
        slide,
        8.8,
        1.4,
        3.8,
        4.95,
        "校正與評估",
        [
            "8 參數 trilinear correction",
            "可表示空間與平滑 residual 誤差界",
            "非連網裝置影響學習",
            "hybrid residual、MAE/RMSE/correlation、IDW、推薦排序",
        ],
        body_size=13,
    )
    add_footer(slide, 21)
    add_formula_walkthrough(prs, 21, compact=False)
    renumber_footers(prs)
    return prs


def build_outline() -> str:
    e8_summary = read_json(DATA / "e8_intervention_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    window_in_domain, window_out_of_domain = window_temperature_domain_counts(window_summary)
    slides = [
        ("封面", ["題目、姓名、雙指導教授、研究定位"]),
        ("研究問題與動機", ["非連網裝置無法直接回報狀態", "有限感測器下仍需估計全室環境", "早期純插值與 local-only 模型都不合理"]),
        ("論文整體邏輯：問題、方法、證據與結論邊界", ["RQ1--RQ3 為主要研究線，RQ4 為次要服務線", "每個研究問題對應方法、E1--E9 證據層與可支持／不可過度宣稱的結論"]),
        ("房間拓樸、感測器與目標區域", ["8 顆角落感測器", "三個主要區域與三個核心裝置"]),
        ("數學模型", ["變數專屬 nominal model", "trilinear correction", "裝置與家具模組化", "溫度、濕度、照度分別使用不同公式"]),
        ("模型學習、推論與推薦資料流", ["學習端：raw records → 對齊 → scenario state → labels → coefficients/checkpoint", "推論端：runtime input → nominal field → correction / hybrid → point or zone prediction", "推薦端：sample / cluster + T/H/L 目標 → 反事實重跑 → penalty reduction 排序"]),
        ("系統實作與介面", ["MCP 是工具化介面，不是預測模型本身", "initialize：設定 scenario、室內 baseline、外部邊界、設備/家具、預設時間與 estimator", "AC state：模式、目標溫度、風量、水平/垂直角度與固定/擺動", "sample point：查指定座標在特定時間或穩定態的溫濕照度", "learn impacts：start/finish before-after record", "window direct / rank actions：輸入外部窗戶資料；rank actions 需指定 sample 與 T/H/L 目標", "Gemma bridge 與 Web demo 分別負責 AI tool calling 與人機展示"]),
        ("learn_impacts：動作如何成為資料記錄", ["start：device_name + device_state 記錄實際操作狀態", "record：儲存 learning_record_id、baseline、外部邊界、家具、elapsed time 與 before observations", "finish：用同一批感測器 after observations 計算 after-before delta", "least squares：由 influence envelope 與 delta 求 learned_device_impacts"]),
        ("驗證流程與比較原則", ["E1-E3：synthetic full-field、IDW baseline、ablation", "E4：非連網裝置影響學習與推薦排序", f"E5：{window_summary.get('count', 0)} 組窗戶矩陣（{window_in_domain} 範圍內／{window_out_of_domain} 範圍外壓力測試）與 direct input", "E6：hybrid residual no-Fourier 與 LOO cross-validation", "E7：bedroom_01 7 天真實快照與 pillow hold-out", f"E8 execution kit：schema / template / analyzer；{e8_summary['trial_counts']['completed']} trials、{e8_summary['evidence_status']}", "E9 public task-aligned benchmark；demo 不是量化實驗"]),
        ("主要結果", ["平均 field MAE", "IDW / Base / LOO Hybrid 誤差比較", "真實臥室 pillow MAE 比較", "推薦排序目前為 counterfactual simulation", "3D 視覺化案例"]),
        ("Hybrid Residual 結果", ["default held-out、no-Fourier、LOO MAE", "train/test sample count", "研究定位不是黑盒替代", "LOO 結果限標準情境 family", "E7 date-block bootstrap 的三因子改善區間下界均大於 0", "E7 逐日剔除的最小 MAE 降幅仍為 T 0.6123、H 3.5551、L 290.5716"]),
        ("公開資料任務拆解", ["SML2010：S1 純照度劣勢、S2 長視窗溫度部分優勢、S3 事件 delta 主要優勢", "Oh2024-inspired transfer：15min 兩點溫度最佳、60min 本研究 readout 最佳、24h persistence 最佳", "次日 primary 與 post-primary adaptive 均未建立優勢；未選中 bias correction 僅約 1% 改善", "RNN 與其他模型共用四筆 history、split、targets、test rows；12/12 parity 通過，RNN lowest MAE 0/12", "CU-BEMS：C1/C3 勝 linear regression 但不勝 persistence，C2 照度劣勢", "明確說明 public benchmark 不是 full 3D 場驗證"]),
        ("研究貢獻與資料策略", ["三因子、有限感測器、非連網裝置、服務化", "canonical synthetic benchmark + real-bedroom snapshots + task-aligned public datasets", "室內應用溫度限 20–30 °C；人體舒適採目標帶與 tolerance", "明確列出每種資料支援的驗證範圍"]),
        ("結論與未來工作", ["長期真實資料、dense real-room ground truth、更多因子、multi-zone、推薦動作介入驗證、閉環控制", "候選動態植物生長情境需補 PPFD/CO2/基質/生物 endpoint", "Kalman family 先定義 state/observation/noise，再做同資料比較"]),
    ]
    slides.extend(
        (
            title,
            [
                left_title,
                right_title,
            ],
        )
        for title, left_title, _, right_title, _ in FORMULA_WALKTHROUGH
    )
    lines = ["# 論文報告投影片大綱", ""]
    for index, (title, bullets) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        lines.extend([f"- {item}" for item in bullets])
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def build_outline_30min() -> str:
    e8_summary = read_json(DATA / "e8_intervention_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    window_in_domain, window_out_of_domain = window_temperature_domain_counts(window_summary)
    slides = [
        ("封面", ["題目、姓名、雙指導教授、研究定位"]),
        ("報告流程", ["背景、文獻、方法、實作、驗證、結論、公式與指標整理"]),
        ("論文整體邏輯：問題、方法、證據與結論邊界", ["研究缺口 → RQ1--RQ4 → 方法核心 → E1--E9 → 有界結論", "controlled、real snapshot、public aligned 與 future intervention 證據層不可互換"]),
        ("研究背景與問題", ["非連網裝置造成空間影響但無法直接讀取", "有限感測器仍需估全室環境"]),
        ("研究問題與貢獻", ["RQ1-RQ4、主要技術貢獻、task-aligned benchmark 策略"]),
        ("文獻定位、研究缺口與比較原則", ["IEQ 實驗、場重建、hybrid model、digital twin 平台之差異", "公開資料集只比較相容子任務"]),
        ("整體系統架構", ["top-down tree 呈現情境觀測、估測學習、服務決策三個責任域", "scripts、Web、MCP/Gemma 共用同一套 estimator path"]),
        ("主要執行資料流", ["runtime request 到 dashboard / MCP response 的流程"]),
        ("房間拓樸、感測器與目標區域", ["8 顆角落感測器與三個區域"]),
        ("模組化裝置與家具阻擋", ["裝置模組化、家具自適應阻擋"]),
        ("數學模型", ["變數專屬 nominal model + residual correction", "早期純插值與 local-only 模型失敗後的調整", "避免把同一套公式套用到溫度、濕度、照度"]),
        ("方法選擇：為什麼不是純插值、純物理或純黑盒", ["IDW 適合作 baseline 但缺設備與方向資訊", "完整 CFD/ray tracing 對低成本即時服務太重", "hybrid residual 只學剩餘誤差，不取代可解釋主模型"]),
        ("模型學習、推論與推薦資料流", ["學習資料流：raw data → 對齊 → scenario state → labels → coefficients/checkpoint", "推論資料流：runtime input → nominal field → correction/hybrid → 溫濕照度", "推薦資料流：sample / cluster + T/H/L 目標 → 反事實重跑 → penalty reduction 排序"]),
        ("系統實作與介面", ["MCP 是工具化介面，不是預測模型本身", "initialize：設定 scenario、baseline、外部邊界、設備/家具、時間與 estimator", "AC state：模式、目標溫度、風量、水平/垂直角度與固定/擺動", "sample point：註冊環境後查指定座標三因子估計", "learn impacts：以 before/after observations 建立可學習資料", "window direct / rank actions：直接輸入窗戶外部資料；rank actions 需指定 sample 與 T/H/L 目標", "Gemma/Ollama 透過 bridge 呼叫 tools；Web demo 負責人機互動展示"]),
        ("learn_impacts：動作如何成為資料記錄", ["start：device_name + device_state 記錄實際操作狀態", "record：儲存 learning_record_id、baseline、外部邊界、家具、elapsed time 與 before observations", "finish：用同一批感測器 after observations 計算 after-before delta", "least squares：由 influence envelope 與 delta 求 learned_device_impacts"]),
        ("驗證設計", ["E1-E3：truth-adjusted simulation、IDW、synthetic ablation", "E4-E6：裝置影響學習、window matrix、hybrid no-Fourier/LOO", "E7：bedroom_01 7 天真實快照與 pillow 位置比較", "E8：推薦動作 before/after intervention protocol", "E9：public datasets 僅作 task-aligned benchmark", "Web demo 與 3D 展示是呈現層，不列為量化實驗"]),
        ("證據鏈與驗證範圍", ["Synthetic full-field 支援完整 3D 場比較，但不等同長期真實場", "Real-bedroom snapshot 支援稀疏校正的 held-out 點位檢查，但不是 dense truth", "Public datasets 僅支援相容子任務，不是單房間 8 點拓樸驗證", "Recommendation 目前是反事實排序，仍需 before/after 介入驗證"]),
        ("情境設計與輸入模式", [f"8 組 scenario、{window_summary.get('count', 0)} 組窗戶矩陣（{window_in_domain} 範圍內／{window_out_of_domain} 範圍外壓力測試）、direct input、timeline"]),
        ("主要量化結果", ["圖表資料：8 組標準情境、full 3D grid Field MAE、log-scale y 軸", "三種柱狀結果：IDW、Base、LOO Hybrid", "真實臥室 raw vs corrected pillow MAE、date-block bootstrap 與逐日剔除敏感度", "推薦有效性以 actual comfort-penalty reduction 驗證", "實驗 E1-E7 與 E9 已有數值輸出；E8 僅為介入 protocol"]),
        ("真實臥室快照與推薦驗證狀態", ["E7：pillow hold-out 不參與 8 角點 fitting；20,000 次 date-block bootstrap 報告三因子 MAE 降幅區間與改善快照數", "E7：7-fold 逐日剔除後，三因子最小 MAE 降幅仍為 0.6123 / 3.5551 / 290.5716", "E7 仍限單一房間、單一 pillow 與七個日期；不是 dense truth 或介入成功率", f"E8：versioned schema、空白 template 與 analyzer 已完成；{e8_summary['trial_counts']['completed']} trials、{e8_summary['evidence_status']}", "真實 before/after 與 matched controls 完成前不得宣稱 efficacy"]),
        ("3D 視覺化結果", ["溫度與照度熱區案例"]),
        ("Hybrid Residual 結果", ["default held-out、no-Fourier、LOO robustness checks", "train/test sample count 與 synthetic benchmark 限制", "LOO 結果限標準情境 family", "真實快照作為 sparse calibration 驗證"]),
        ("公開資料任務拆解：SML2010", ["原 E9：S1 照度弱、S2 混合、S3 event delta 最強", "Oh2024-inspired transfer：15min 兩點溫度最低 MAE", "60min 由本研究 readout 最佳；24h 由 persistence 最佳且 transfer 劣於 raw physics", "次日 primary 選中 trend 但 test 惡化 7.34% / 8.36%，bootstrap interval 均跨 0", "RNN 與其他模型共用四筆 history、split、targets、test rows；12/12 parity 通過，RNN lowest MAE 0/12", "資料 confidential；方法移植不等於原文 CNN--LSTM 重現"]),
        ("公開資料任務拆解：CU-BEMS", ["C1：AC 溫濕度可補強 linear regression", "C2：商辦照度與單房間假設差距大", "C3：compound event 可勝 linear regression 但不勝 persistence"]),
        ("結論、限制與未來工作", ["目前完成度、真實快照限制、hybrid 泛化限制、推薦動作尚需介入驗證、task-aligned benchmark 與後續方向", "室內溫度限 20–30 °C；人體舒適採 tolerance，RNN 負向結果保留", "候選植物生長情境需補 PPFD/CO2/基質/生物 endpoint；Kalman 尚未評估"]),
        ("公式與指標整理", ["場模型：三因子場、總估計式、baseline、activation、envelope", "三因子公式：溫度、濕度、照度分別說明", "校正與評估：8 點三線性校正、影響學習、hybrid residual、metrics、IDW、推薦排序"]),
    ]
    slides.extend(
        (
            title,
            [
                left_title,
                right_title,
            ],
        )
        for title, left_title, _, right_title, _ in FORMULA_WALKTHROUGH
    )
    lines = ["# 論文報告投影片大綱（30min 版）", ""]
    for index, (title, bullets) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        lines.extend([f"- {item}" for item in bullets])
        lines.append("")
    return "\n".join(lines).strip() + "\n"


SPEAKER_NOTE_GLOSSARY: List[Tuple[str, Tuple[str, ...], str]] = [
    ("單房間", ("單房間", "single-room"), "本研究限定在單一矩形房間，不處理多房間或整棟建築的氣流與能量交換。"),
    ("非連網家電/裝置", ("非連網家電", "非連網裝置", "非連網"), "沒有穩定 API 或遙測資料可讀取狀態的冷氣、窗戶、照明等設備。"),
    ("稀疏感測", ("稀疏感測", "少量感測", "sparse"), "感測點數量少於完整空間場需求，需靠模型與校正推估未量測位置。"),
    ("角落感測器", ("角落感測器", "8 顆", "8 點", "8個角點", "8 個角點"), "配置在房間地面四角與天花板四角的 8 個感測點，用於建立 sparse observation 與 residual correction。"),
    ("空間數位孿生", ("空間數位孿生", "數位孿生", "digital twin"), "以房間幾何、裝置、感測器與模型維持一個可查詢的室內環境狀態估計。"),
    ("三因子", ("三因子", "temperature", "humidity", "illuminance", "溫度、濕度、照度", "溫濕照度"), "本研究同時估計溫度、相對濕度與照度三種室內環境量。"),
    ("空間場", ("空間場", "field", "場估計", "場重建", "field MAE"), "不是單一平均值，而是在房間 3D 座標中任意位置可查詢的環境量分布。"),
    ("控制導向", ("控制導向", "control-oriented"), "模型重點在支援查詢、比較與推薦排序，而不是取代高精度物理模擬器。"),
    ("CFD", ("CFD",), "Computational Fluid Dynamics，計算流體力學；可模擬細緻氣流，但邊界條件與計算成本高。"),
    ("Ray tracing", ("ray tracing", "光線追蹤"), "依光線路徑追蹤照明傳播的精密光學方法；本研究只採輕量照度幾何與一次漫反射近似。"),
    ("API", ("API",), "Application Programming Interface，讓系統讀取或控制設備狀態的程式介面。"),
    ("遙測", ("遙測", "telemetry"), "設備主動回報狀態或感測資料；非連網裝置通常缺少這類資料。"),
    ("Sample", ("sample point", "sample ", "sample/cluster", "sample 或 zone"), "指定房間中的查詢點，用來取得該座標的三因子估計。"),
    ("Zone", ("zone", "區域", "目標區域"), "房間內的目標區域，用於彙整多個點的平均狀態或舒適度評估。"),
    ("Baseline", ("baseline", "Indoor baseline", "室內 baseline"), "泛指比較或模型參考基準；在模型脈絡中常指未加入設備作用前的室內溫濕照度基準。"),
    ("外部邊界", ("外部邊界", "外部環境邊界", "outdoor"), "室外溫度、濕度、日照等會透過窗戶或邊界條件影響室內的輸入。"),
    ("Nominal model", ("nominal model", "N_v", "Nᵥ"), "主模型的可解釋估計部分，描述設備、邊界與空間結構造成的主要趨勢。"),
    ("Residual", ("residual", "剩餘誤差", "殘差"), "觀測或 truth 與模型預測之間的差，用於校正或第二層學習。"),
    ("Residual correction", ("residual correction", "校正場", "感測器校正"), "利用感測器 residual 修正 nominal model，使估計更貼近觀測。"),
    ("Trilinear correction", ("trilinear", "三線性", "三線性補間", "三線性校正"), "用 X/Y/Z 三個座標方向的一階補間，由 8 個角點 residual 推估室內 residual 場。"),
    ("Power calibration", ("power calibration", "power scale", "P_ac", "P_win", "P_light"), "依觀測差異調整設備影響強度，避免裝置作用尺度只依預設值決定。"),
    ("AC operating state", ("AC state", "ac_mode", "target_temperature", "fan_speed", "fan_strength", "出風角度", "擺動"), "冷氣操作狀態包含模式、設定溫度、風速/風量、水平與垂直出風角度，以及 fixed/swing 擺動設定。"),
    ("device_state", ("device_state",), "learn_impacts start 階段輸入的裝置操作狀態，包含 activation、kind、power 與冷氣模式、設定溫度、風速、風向等欄位。"),
    ("device_specs", ("device_specs",), "系統把 device_state 合併進目前註冊設備後形成的完整設備清單，是後續 sample、learning 與 ranking 使用的 runtime 裝置狀態。"),
    ("learning_record_id", ("learning_record_id",), "每一次 learn_impacts start 產生的唯一紀錄編號，用來在 finish 階段把 after observations 接回同一筆事件。"),
    ("before_observations / after_observations", ("before_observations", "after_observations"), "同一批感測器在裝置操作前後的真實讀值，格式通常是 sensor name 對應 temperature、humidity 與 illuminance。"),
    ("learned_device_impacts", ("learned_device_impacts", "metric coefficients"), "由 before/after 差值與設備 influence envelope 解出的裝置影響係數，描述該操作對三因子的方向與大小。"),
    ("After-before delta Δy", ("Δy", "after-before delta"), "裝置操作後觀測值減去操作前觀測值，用來把實際動作轉成可學習的數值變化。"),
    ("設計矩陣 X", ("X_{i,k}", "設計矩陣", "X β_m", "Xβ"), "每一列對應一個感測點或樣本，每一欄對應一個裝置 influence envelope，用於 least-squares 估計。"),
    ("影響係數 β_m", ("β_m", "argmin_{β}", "影響係數"), "描述裝置操作對第 m 個環境因子的方向與大小；m 可為溫度、濕度或照度。"),
    ("Hybrid residual", ("hybrid residual", "F_hybrid", "hybrid"), "在可解釋 base estimator 後面再加一個資料驅動 residual 模型，不直接取代主模型。"),
    ("MCP", ("MCP", "Model Context Protocol"), "Model Context Protocol，是讓 LLM application 以標準化方式連接外部資料與工具的 open protocol；本研究用它封裝數位孿生工具。"),
    ("MCP host/client/server", ("MCP host", "MCP client", "MCP server", "client-server", "server"), "MCP 採 client-server 概念；host/client 是使用工具的 AI 應用端，server 則暴露工具、資源或 prompt。"),
    ("MCP Tools", ("MCP tools", "tools/list", "tools/call", "initialize_environment", "sample_point", "learn_impacts", "rank_actions"), "MCP server 可暴露可執行工具；client 可列出工具並以結構化 arguments 呼叫。"),
    ("MCP Resources", ("resources/list", "Resources", "resources"), "MCP resources 是 server 提供給 client 的上下文資料，例如檔案、資料庫 schema 或應用資料。"),
    ("MCP Prompts", ("prompts/list", "Prompts", "prompt templates"), "MCP prompts 是 server 提供的結構化 prompt template，可由 client 取得並填入參數。"),
    ("JSON-RPC", ("JSON-RPC", "jsonrpc"), "MCP 使用 JSON-RPC 2.0 編碼 request、response 與 notification。"),
    ("stdio transport", ("stdio", "stdin", "stdout", "standard input", "standard output"), "MCP 的本地 transport 之一；client 啟動 server subprocess，透過 stdin/stdout 傳送 UTF-8 JSON-RPC 訊息。"),
    ("Streamable HTTP", ("Streamable HTTP", "HTTP MCP", "遠端 HTTP MCP"), "MCP 的另一種標準 transport，適合遠端或網路化部署。"),
    ("Protocol version", ("protocolVersion", "2024-11-05", "2025-06-18"), "MCP 會在 initialize 階段協商 protocolVersion；本研究本地 server 目前回傳既有版本並實作 tools workflow。"),
    ("Web demo", ("Web demo", "dashboard"), "人機互動展示介面，用於查看 3D 場、時間軸、設備狀態與查詢結果。"),
    ("Gemma/Ollama bridge", ("Gemma", "Ollama", "bridge"), "讓本地語言模型透過工具呼叫流程存取模型服務的橋接層。"),
    ("Tool calling", ("tool calling", "AI 工具呼叫", "AI agent"), "語言模型不是直接計算答案，而是呼叫外部工具取得模型查詢或操作結果。"),
    ("服務編排", ("服務編排", "service", "service orchestration"), "把 scenario、模型估計、校正、推薦與輸出流程串接起來的中介層。"),
    ("Estimator", ("estimator",), "實際負責產生場估計的模型物件，可切換 base、corrected 或 hybrid 版本。"),
    ("Scenario", ("scenario", "標準情境", "情境"), "一組房間、設備、外部邊界與時間設定，用於模擬或驗證。"),
    ("Direct input", ("direct input", "window direct"), "不使用預設矩陣情境，直接輸入外部溫濕度、日照與開窗比例。"),
    ("Window matrix", ("window matrix", "窗戶矩陣"), "依季節、天氣、時段等組合建立的窗戶外部邊界情境集合。"),
    ("Quasi-steady state", ("quasi-steady", "穩定態"), "近似達到穩定但非嚴格物理穩態的狀態，用於簡化時間響應解讀。"),
    ("Activation", ("activation", "Aⱼ", "A_j"), "設備啟動強度隨時間接近穩態的函數。"),
    ("Influence envelope", ("influence envelope", "Eⱼ", "E_j", "envelope"), "設備對某位置的空間作用權重，通常含時間強度、距離、方向與遮蔽。"),
    ("Distance decay", ("距離衰減", "Rⱼ", "R_j"), "距離設備越遠，局部作用越弱的權重函數。"),
    ("Directionality", ("方向性", "Dⱼ", "D_j"), "冷氣出風方向、窗戶日照方向或光源方向造成的非均向影響。"),
    ("Visibility/obstruction", ("可見性", "遮蔽", "阻擋", "Vⱼ", "V_j", "obstruction"), "家具或幾何遮擋造成設備影響變弱的因素。"),
    ("燈具光束權重", ("Φ_light", "Q_light", "V_light", "光束角", "cosine"), "照度模型中用方向權重、距離衰減與可見性近似燈具直射光，而非完整光線追蹤。"),
    ("IDW", ("IDW",), "Inverse Distance Weighting，反距離加權插值；只使用距離與感測值，不含設備物理先驗。"),
    ("黑盒模型", ("黑盒", "純黑盒"), "主要依資料學習輸入輸出關係、但內部物理意義較不明確的模型。"),
    ("IEQ", ("IEQ",), "Indoor Environmental Quality，室內環境品質，通常涵蓋熱舒適、空氣品質、照明等因素。"),
    ("Hybrid thermal model", ("hybrid thermal model",), "結合物理結構與資料驅動方法的熱環境模型。"),
    ("Task-aligned benchmark", ("task-aligned benchmark", "相容子任務", "task-aligned"), "只比較公開資料集中與本研究觀測型態相容的子任務，不宣稱完整場驗證。"),
    ("SML2010", ("SML2010",), "公開智慧建築資料集；本研究用於 two-point boundary-response 類任務。"),
    ("CU-BEMS", ("CU-BEMS",), "商辦建築能源管理資料集；本研究用於 zone-level device-response 類任務。"),
    ("Public dataset", ("Public datasets", "公開資料", "公開資料集"), "外部公開資料來源；本研究只用於相容任務壓力測試，不作完整 3D truth。"),
    ("Synthetic full-field", ("Synthetic full-field", "truth-adjusted", "受控完整場", "synthetic"), "可取得完整場 truth 的受控驗證資料，用於完整 3D 場誤差比較。"),
    ("Ablation", ("ablation", "消融"), "移除或替換模型元件後比較性能，用來檢查各元件貢獻。"),
    ("LOO", ("LOO", "leave-one-scenario-out"), "Leave-one-scenario-out，輪流留下一個情境作測試，檢查模型是否只對單一切分有效。"),
    ("No-Fourier", ("No-Fourier", "Fourier"), "去除或比較 Fourier 相關處理的對照設定，用於確認改善不是單一頻域技巧造成。"),
    ("Real-bedroom snapshot", ("Real-bedroom snapshot", "bedroom_01", "真實臥室", "快照"), "真實臥室中的稀疏量測快照，用於檢查校正對未參與 fitting 點位的改善。"),
    ("Pillow hold-out", ("pillow", "hold-out"), "將 pillow 位置作為未參與校正 fitting 的參考點，用於測試非感測點估計效果。"),
    ("Dense ground truth", ("dense", "dense truth", "dense real-room ground truth", "完整 3D 場真值"), "房間內大量點位的真實環境場資料，是更嚴格但較難取得的驗證基準。"),
    ("Log-scale", ("log-scale",), "對數刻度；適合把量級差很多的誤差放在同一張圖，但柱高不能用線性比例直接比較。"),
    ("MAE", ("MAE",), "Mean Absolute Error，平均絕對誤差；數值越低代表平均偏差越小。"),
    ("RMSE", ("RMSE",), "Root Mean Squared Error，均方根誤差；比 MAE 更放大尖峰或離群誤差。"),
    ("Correlation", ("Correlation", "Corr=", "cov(ŷ,y)", "相關係數"), "衡量預測與真值趨勢方向一致性的指標。"),
    ("Persistence", ("persistence",), "直接沿用上一時步值作預測的時間序列 baseline，在高慣性短視窗資料中通常很強。"),
    ("Linear regression", ("linear regression", "LR"), "線性回歸 baseline，用線性權重將輸入特徵映射到目標值。"),
    ("Structured prior", ("structured prior",), "模型內建的設備、邊界與物理結構先驗。"),
    ("Facade event delta", ("facade event delta", "event delta", "delta response"), "檢查外部邊界或事件造成的變化量，而非只預測下一時步絕對值。"),
    ("Device-response benchmark", ("device-response", "zone-level device-response"), "以設備用電或啟動訊號對 zone 環境變化的響應作為比較任務。"),
    ("Zone-level", ("zone-level",), "以建築區域平均值為資料粒度，不含房間內細緻 3D 幾何。"),
    ("Counterfactual simulation", ("counterfactual", "反事實"), "假設某候選動作發生後重新估計結果，用來比較預期改善。"),
    ("Comfort penalty", ("comfort penalty", "Penalty", "舒適度"), "偏離目標溫濕照度時的懲罰值，推薦排序用它衡量改善幅度。"),
    ("區域平均 q_m(S)", ("q_m", "q_base", "q_a", "sample scope"), "在指定 sample scope 或目標區域 S 內，彙整第 m 個環境因子的平均狀態。"),
    ("目標值與容許範圍", ("g_m", "δ_m", "容許範圍"), "g_m 是舒適目標，δ_m 是允許偏離範圍；超出範圍才累積 penalty。"),
    ("Before/after intervention", ("before/after", "intervention", "介入驗證"), "實際採取動作前後量測環境變化，用於驗證推薦是否有因果改善效果。"),
    ("ESP32", ("ESP32",), "低成本微控制器平台，可用於後續長期真實感測資料蒐集。"),
    ("CO2", ("CO2", "CO₂"), "二氧化碳濃度，可作為未來室內空氣品質因子。"),
    ("PM2.5", ("PM2.5",), "細懸浮微粒濃度，可作為未來室內空氣品質因子。"),
    ("Multi-zone model", ("multi-zone", "partition model"), "將房間或建築切成多個區域處理交換與隔間效應的模型。"),
    ("閉環控制", ("閉環控制", "closed-loop"), "模型輸出進一步驅動控制動作，並用後續感測結果回饋修正決策。"),
    ("座標系", ("座標系", "原點", "p = (x,y,z)", "p=(x,y,z)"), "用 x/y/z 公尺座標描述房間內位置；本研究原點在地面西南角。"),
    ("T(p,t)", ("T(p,t)", "T："), "位置 p、時間 t 的溫度場值。"),
    ("H(p,t)", ("H(p,t)", "H："), "位置 p、時間 t 的相對濕度場值。"),
    ("L(p,t)", ("L(p,t)", "L："), "位置 p、時間 t 的照度場值。"),
    ("b₀", ("b₀", "T₀", "H₀", "L₀"), "室內基準狀態，包含基準溫度、相對濕度與照度。"),
    ("ζ", ("ζ",), "高度正規化座標，用於描述查詢點相對於房間高度的位置。"),
    ("τⱼ", ("τⱼ", "τ_j"), "設備時間響應常數，控制 activation 接近穩態的速度。"),
    ("B 項", ("B_T", "B_H", "B_ac", "B_win", "B_light", "全室項"), "表示全室平均狀態偏移的 bulk/global effect。"),
    ("S 項", ("S_T", "S_H", "S_ac", "S_win", "S_light", "局部項"), "表示設備附近、窗邊或光源附近的局部空間差異。"),
    ("clip[0,100]", ("clip[0,100]",), "把相對濕度限制在 0% 到 100% 的合理物理範圍內。"),
    ("max{0}", ("max{0", "max{0,"), "把照度限制為非負值，避免模型輸出不合理的負照度。"),
    ("直射光", ("直射光", "direct source", "L_winᵈ", "L_lightᵈ"), "直接由窗戶或燈具到達查詢點的照度貢獻。"),
    ("環境光", ("環境光", "ambient"), "非單一路徑直射、較均勻分布的背景照度貢獻。"),
    ("一次漫反射", ("一次漫反射", "Iʳ", "reflect", "漫反射"), "只計算一次表面反射對照度的回填效果，是輕量近似而非完整 radiosity。"),
    ("反射率 ρ", ("ρ_s", "ρ"), "表面反射率，描述牆面、地板或家具把入射光反射出去的比例。"),
    ("三線性函數空間", ("𝒱", "span{1", "XYZ"), "由 1、X、Y、Z 與交互項組成的 8 維 residual 表示空間。"),
    ("角點 residual", ("rᵛ", "r^", "rᵛ_", "p_{abc}"), "角落感測器觀測值與 nominal model 預測值的差。"),
    ("補間權重 ℓ", ("ℓ₀", "ℓ₁"), "三線性補間中每個角點 residual 對內部點的權重函數。"),
    ("誤差上界", ("誤差上界", "M_xx", "M_yy", "M_zz"), "用 residual 二階曲率限制三線性補間與真實 residual 的最大偏差。"),
    ("特徵向量 φᵢ", ("φᵢ", "feature vector"), "模型訓練時輸入的情境特徵，例如座標、時間、baseline、外部條件與設備作用。"),
    ("標籤 yᵢ", ("yᵢ", "label", "labels"), "監督式學習中的目標值；本研究常用 true field 與 base estimator 的差作 residual label。"),
    ("損失函數 ℒ", ("ℒ", "loss"), "訓練 neural network 時要最小化的目標函數。"),
    ("正則化 λ", ("λ", "regularization", "正則化"), "限制模型參數大小以降低過擬合的項。"),
    ("IDW 權重", ("w_s", "q：距離權重"), "IDW 中由距離決定的感測器權重，距離越近權重越高。"),
    ("Score(a)", ("Score(a)", "Score"), "候選動作 a 的推薦分數，通常由採取前後 comfort penalty 的下降量決定。"),
]


def glossary_notes_for_slide(title: str, paragraphs: Sequence[str]) -> List[Tuple[str, str]]:
    text = "\n".join([title, *paragraphs]).lower()
    notes: List[Tuple[str, str]] = []
    for label, aliases, explanation in SPEAKER_NOTE_GLOSSARY:
        if any(alias.lower() in text for alias in aliases):
            notes.append((label, explanation))
    return notes


def build_speaker_notes_30min() -> str:
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    window_in_domain, window_out_of_domain = window_temperature_domain_counts(window_summary)
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    e8_summary = read_json(DATA / "e8_intervention_summary.json")
    avg_mae = average_field_mae(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]
    bedroom_bootstrap = bedroom_aggregate["paired_day_block_bootstrap"]
    bedroom_bootstrap_metrics = bedroom_bootstrap["metrics"]
    bedroom_lodo = bedroom_aggregate["leave_one_date_out_sensitivity"]
    bedroom_lodo_metrics = bedroom_lodo["metrics"]
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]
    scenario_count = submission_summary["base_ablation"]["scenario_count"]
    base_variants = submission_summary["base_ablation"]["variants"]
    idw_mae = base_variants["idw"]["average_field_mae"]
    base_mae = base_variants["full_base"]["average_field_mae"]
    loo_hybrid_mae = loo["average_hybrid_field_mae"]
    grid_resolution = bedroom_summary["grid_resolution"]

    slides: List[Tuple[str, List[str]]] = [
        (
            "封面",
            [
                "各位老師好，我是林昀佑。今天報告的題目是「單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型」。題目中的「單房間」代表研究範圍先限制在一個可明確定義幾何邊界的室內空間；「非連網家電」指的是冷氣、窗戶、燈具這類會影響環境，但通常沒有穩定 API 可以讀取狀態的設備；「稀疏感測」則表示系統只依賴少量感測點，而不是完整佈滿房間的感測陣列。",
                "這個研究要解決的核心問題是：使用者真正關心的是床邊、桌面、窗邊或其他位置的舒適狀態，但系統通常只能看到少數感測器讀值，而且很多設備狀態還要靠使用者或介面輸入。因此我建立一個可以把房間幾何、設備狀態與少量感測資料結合起來的空間數位孿生原型。",
                "整體研究不是要取代完整 CFD 或精密光學模擬，也不是把所有問題交給黑盒模型處理。我的定位是控制導向與決策支援：模型需要足夠可解釋、可以被感測資料校正，也能輸出任意位置或區域的溫度、濕度與照度估計。",
                "後面報告會先說明為什麼這個問題重要，再說明模型與系統怎麼設計，最後用 synthetic full-field、真實臥室快照與公開資料相容子任務說明目前的驗證範圍。",
            ],
        ),
        (
            "報告流程",
            [
                "整份報告會依照「問題、方法、實作、驗證、限制」的順序走。前半段先建立研究動機，說明為什麼非連網設備與有限感測器會讓室內環境估計變困難。",
                "接著會進入文獻定位與研究缺口。這裡的重點不是逐篇列文獻，而是說明既有 IEQ、場重建、hybrid model 與 digital twin 研究，和本研究的單房間、三因子、低成本角落感測設定有什麼差異。",
                "方法部分會先看整體系統架構，再拆成數學模型、校正流程、影響學習與推薦排序。這樣安排是因為本研究不是只有一個公式，而是一個從輸入資料到可查詢服務的完整 pipeline。",
                "驗證部分會分成受控完整場、真實臥室快照、hybrid residual、公開資料集與推薦驗證狀態。每一類資料能支持的結論不同，所以我會特別區分哪些結果可以證明 field reconstruction，哪些只能作為 sparse calibration 或 task-aligned benchmark。",
                "最後會整理目前已完成的貢獻、尚未完成的真實介入驗證，以及後續要補強的硬體、資料量與泛化能力。",
            ],
        ),
        (
            "論文整體邏輯：問題、方法、證據與結論邊界",
            [
                "這張圖是整篇論文的 argument map。起點不是 MCP 或 Web，而是一般房間同時存在稀疏感測與非連網裝置兩個限制；因此需要回答空間場估計、裝置影響學習與決策支援三個主要研究問題。",
                "RQ1 對應變數專屬 nominal model、power calibration、trilinear correction 與 optional hybrid residual，證據來自 E1 到 E3、E6 與 E7。這些結果可以支持受控完整場與真實未見點改善，但不能外推成任意房間 dense truth。",
                "RQ2 對應 before/after delta 與裝置 spatial basis，證據包含 E4、E5 與 E9 的事件型相容子任務。這可以支持 structured impact prior 的價值，但不等同已完成真實因果識別。",
                "RQ3 對應 point 或 zone sample、完整 T/H/L target 與反事實 comfort-penalty 排序。目前只能說模型可提供決策支援；E8 before/after intervention 尚未完成，所以不能宣稱推薦具有實際因果改善。",
                "RQ4 是 secondary systems line，說明 scripts、Web、MCP 與 Gemma bridge 如何共用同一 estimator。它證明介面整合與模型重用，但不是論文的 headline novelty。",
            ],
        ),
        (
            "研究背景與問題",
            [
                "一般智慧居家或智慧建築如果要做舒適度評估、能源管理或設備控制，必須先知道室內環境狀態。但實際房間裡的很多設備並不會主動回報狀態，例如傳統冷氣、手動窗戶或一般燈具。",
                "這造成第一個問題：設備有影響，但系統看不到完整狀態。冷氣可能正在冷房、窗戶可能正在引入外氣或日照、燈具可能改變桌面照度，但如果沒有 API 或遙測，就不能直接把這些資訊當成可靠輸入。",
                "第二個問題是感測器很少。即使有 1 到 8 顆感測器，它們也只能代表少數位置。以臥室為例，床頭、書桌、窗邊和門側的狀態可能不同，但使用者通常不會在每個位置都放一顆感測器。",
                "所以本研究要處理的是一個資訊不完整的空間估計問題：在設備狀態需要被描述或推定、感測點又有限的情況下，如何估計整個房間的三因子分布。",
                "這裡的目標不是只預測單一平均溫度，而是要能回答「某個位置現在大概是多少」、「某個區域是否偏離目標」、「如果調整冷氣或窗戶，哪個動作比較可能改善」這類空間化問題。",
            ],
        ),
        (
            "研究問題與貢獻",
            [
                "本研究拆成四個研究問題。第一是場估計問題：在單房間中，只使用 8 顆角落感測器與房間幾何，能不能估計溫度、濕度與照度的 3D 空間場。",
                "第二是裝置影響學習問題：對於沒有 API 的冷氣、窗戶與燈具，能不能透過 before/after observations，把實際操作轉成可學習的影響係數，而不是只依賴人工設定。",
                "第三是決策支援問題：當使用者指定 sample point 或 zone，以及三因子的目標值時，系統能不能用反事實模擬比較候選動作，並依 comfort penalty reduction 排序。",
                "第四是系統封裝問題：這個模型能不能不是只停留在離線程式，而是封裝成 Web demo 與 MCP tools，讓人機介面或 AI client 都能查詢同一個核心服務。",
                "對應的貢獻包含：變數專屬三因子 nominal model、8 點 residual correction、非連網裝置影響學習、hybrid residual 修正，以及公開資料集上的 task-aligned benchmark。",
                "同時也要先界定範圍：本研究不宣稱已完成多房間建築級模型，也不宣稱公開資料集能驗證完整 3D 場；公開資料只用來測相容子任務，完整場驗證主要來自受控 synthetic full-field。",
            ],
        ),
        (
            "文獻定位、研究缺口與比較原則",
            [
                "文獻可以先分成幾類。IEQ 研究通常關心熱舒適、空氣品質與照明品質；有限感測器場重建研究關心如何從少量點推估空間分布；hybrid thermal model 結合物理結構與資料驅動；digital twin 平台則強調系統整合與可視化。",
                "這些研究各自提供重要基礎，但和我的設定仍有落差。有些研究依賴比較完整的 BMS 或設備遙測，有些只處理溫度或照度單一因子，有些不處理房間內任意位置的 3D 場估計。",
                "本研究的缺口定位是把幾個限制同時放在一起：單房間、低成本角落感測、非連網裝置、三因子環境量，以及可被查詢和推薦使用的控制導向模型。",
                "公開資料集比較也要小心。SML2010 和 CU-BEMS 很有價值，但它們不是為本研究的 8 點單房間拓樸設計，也沒有 dense 3D ground truth。因此我不把它們說成完整替代驗證，而是拆成 task-aligned 子任務。",
                "這樣做的原因是避免過度宣稱。能直接比較的地方就比較，例如 boundary response 或 device-response；不能比較的地方就明確說明限制，維持論文結論和資料支撐一致。",
            ],
        ),
        (
            "整體系統架構",
            [
                "這頁說明系統不是單一模型函式，而是一棵 top-down abstraction tree。最上層是單房間三因子空間數位孿生系統，往下分成情境與觀測、估測與學習、服務與決策三個責任域。",
                "情境與觀測層負責把 room schema、zones、furniture blockers、8 點角落感測、外部邊界與時間整理成系統可用的狀態。這一層界定資料從哪裡來，也界定本研究的 sparse IoT sensing 前提。",
                "估測與學習層才是主方法核心，包含 T/H/L nominal field model、設備影響函數、active-device power calibration、trilinear correction、非連網裝置影響學習與 optional hybrid residual。",
                "服務與決策層把同一套 estimator path 暴露給 reproduction scripts、Web demo、MCP tools 與 Gemma bridge。這些介面可以使用模型，但不取代模型本身。",
                "這樣畫成樹狀圖的好處是先釐清責任邊界，再到下一頁用 runtime flow 說明執行順序。未來如果要新增 CO2、PM2.5 或新設備類型，也能知道要先補輸入 schema、模型層，還是服務輸出層。",
            ],
        ),
        (
            "主要執行資料流",
            [
                "執行資料流可以分成四步。第一步是取得 runtime request，來源可能是 Web dashboard，也可能是 MCP tool call。兩者都會提供或引用一個 scenario，包含房間、設備、外部邊界與時間。",
                "第二步是資料正規化。系統會把 baseline、device state、furniture obstruction、direct input 和 sensor observations 整理成 estimator 可讀的狀態，而不是讓每個入口各自解讀。",
                "第三步是估計流程。Estimator 先建立 nominal field，這是由設備、幾何與物理先驗產生的主要趨勢；接著用角落感測 residual 做 trilinear correction，必要時再加上 hybrid residual。",
                "第四步是輸出。Dashboard 會拿到 3D 場與視覺化資料，sample_point 會回傳指定座標的 T/H/L，zone summary 會回傳區域平均，rank_actions 則回傳候選動作的預期改善排序。",
                "這頁要強調的是一致性：Web 和 MCP 只是不同使用方式，不是兩套模型。因此同一個 scenario 下，畫面展示、工具查詢和推薦排序應該對齊。",
            ],
        ),
        (
            "房間拓樸、感測器與目標區域",
            [
                "本研究使用單一矩形房間作為主要研究場景，尺寸為 6 m × 4 m × 3 m。座標系以公尺表示，原點設在房間地面西南角，x、y、z 分別對應房間的水平與高度方向。",
                "使用座標系的原因是模型需要回答任意位置的查詢，而不是只回答感測器所在位置。例如桌面、床邊或窗邊都可以被表示成 p=(x,y,z)，再查詢該點的三因子估計。",
                "8 顆感測器放在地面四角與天花板四角。這樣的配置讓模型至少能觀察房間邊界的低處與高處 residual，對 3D 空間補間比只放同一平面更合理。",
                "這 8 點不等於完整量到全室。它們的角色是提供 sparse observation，用來修正 nominal model 的 residual。Nominal model 提供主要趨勢，角點 residual 則用 trilinear correction 補足模型偏差。",
                "目標區域分成窗邊、中心與門側等 zone，目的是讓後續推薦不只看單一點，也能看一個區域的平均舒適度。這對實際使用比較合理，因為使用者通常關心床區、工作區或窗邊區域，而不是單一座標點。",
            ],
        ),
        (
            "模組化裝置與家具阻擋",
            [
                "這頁說明為什麼要把設備和家具都放進模型。冷氣、窗戶與燈具不只是 on/off 狀態，它們在房間中有位置、方向、作用距離、啟動強度和時間響應。",
                "以冷氣為例，除了模式之外，目標溫度、風量、水平與垂直出風角度，以及是否固定或擺動，都會影響哪個位置先變冷。窗戶則和開窗比例、外部溫濕度、日照方向有關；燈具則和光源位置、光束方向與距離衰減有關。",
                "家具在這裡被建模成 bounding box obstruction。它不需要做到精密流體或光學遮擋，但至少可以讓模型知道某些位置和設備之間可能被床、桌子或櫃子阻擋。",
                "因此每個裝置的影響可以被拆成 activation、distance decay、directionality 與 visibility/obstruction 幾個權重。這些權重再組合成 influence envelope，描述裝置對空間中某一點的作用強弱。",
                "這樣做的目的不是追求 CFD 或 ray tracing 等級的細節，而是在低成本資料下保留最重要的幾何資訊，避免模型退化成只看距離或全室平均的粗略估計。",
            ],
        ),
        (
            "數學模型",
            [
                "核心估計式是 F_hat_v(p,t)=N_v(p,t)+C_v(p,t)。v 代表環境變數，可以是 temperature、humidity 或 illuminance；p 是房間內位置；t 是時間或 elapsed time。",
                "N_v 是 nominal model，也就是模型在沒有感測器校正前，根據房間、設備、外部邊界和時間所估出的主要趨勢。它負責把物理與幾何先驗放進估計中，例如冷氣方向、窗戶日照或燈具距離。",
                "C_v 是 correction field，也就是用感測器 residual 建立的校正項。Residual 是觀測值減掉 nominal prediction，如果角落感測器發現模型偏高或偏低，C_v 會把這個偏差用三線性補間延伸到房間內部。",
                "三個因子的 nominal model 不能共用同一套公式。溫度主要處理熱交換、熱源與垂直分層；濕度要處理除濕與外氣水氣交換；照度則要處理光源幾何、遮蔽和一次漫反射。",
                "所以這頁的重點是兩層設計：第一層是變數專屬的可解釋 base estimator，第二層是由稀疏感測器提供的 residual correction。後面的 hybrid residual 則是在這兩層之後再學剩餘誤差。",
            ],
        ),
        (
            "方法選擇：為什麼不是純插值、純物理或純黑盒",
            [
                "這頁是在說明方法選擇。第一個可能方法是純插值，例如 IDW。IDW 的優點是簡單，也適合作為 baseline；但它只知道感測器位置與讀值，不知道冷氣出風、窗戶日照、燈具位置或家具遮蔽。",
                "因此在設備造成局部影響時，純插值會吃虧。例如窗邊強光或冷氣出風口附近的局部冷區，不一定能從距離最近的角落感測器直接推回來。",
                "第二個可能方法是完整物理模擬，例如 CFD 或 ray tracing。這類方法精度潛力高，但需要材料、邊界條件、氣流、反射率等大量資訊，計算成本也比較高，和本研究的低成本即時服務目標不一致。",
                "第三個可能方法是純黑盒模型。問題是本研究目前資料量有限，而且需要解釋設備與空間結構如何影響結果；如果完全黑盒，口徑上比較難說明為什麼某個動作會被推薦。",
                "因此本研究採取折衷：先用可解釋 base model 表達主要物理與幾何趨勢，再用 residual correction 貼近感測器，最後用 hybrid residual 學 base model 尚未捕捉到的剩餘誤差。",
            ],
        ),
        (
            "模型學習、推論與推薦資料流",
            [
                "這頁把系統分成三條資料流：學習、推論與推薦。三條資料流共用同一個房間與設備表示，但目的不同。",
                "學習資料流從 raw records 開始。系統會把 before/after observations、device_state、baseline、外部邊界和時間對齊成 scenario state，再產生訓練 labels。對 learn_impacts 來說，label 是裝置操作前後的感測變化；對 hybrid residual 來說，label 是 truth field 和 base estimator 之間的 residual。",
                "推論資料流處理的是使用者當下查詢。輸入 runtime state 後，模型先產生 nominal field，再套用 correction 或 hybrid residual，最後得到某個 sample point、zone 或整個 3D grid 的三因子估計。",
                "推薦資料流則是建立在推論之上。系統會列出候選動作，例如調整冷氣、開窗或改變照明，對每個候選動作重新跑一次 counterfactual simulation，再計算 comfort penalty 是否下降。",
                "因此推薦結果目前是模型反事實排序，不等於系統已經真的控制設備並量到改善。這個區分很重要，因為後面 E8 會把實際 before/after intervention 列為未來需要完成的驗證。",
            ],
        ),
        (
            "系統實作與介面",
            [
                "MCP 的全名是 Model Context Protocol。它不是我的預測模型，也不是一個新的神經網路架構，而是一個讓 LLM application 用標準化方式連接外部資料與工具的 open protocol。",
                "如果老師問 std 或 standard，我會回答：MCP 本身是標準化的 protocol；官方規格用 JSON-RPC 2.0 表示 request、response 與 notification。它的標準 transport 包含 stdio 和 Streamable HTTP。",
                "stdio 是 standard input/output 的意思，適合本機工具。client 會啟動 MCP server subprocess，server 從 stdin 讀 newline-delimited JSON-RPC message，再把 response 寫到 stdout；stderr 只用於 log。",
                "在我的系統裡，數位孿生核心服務被包成本地 MCP server，主要暴露 tools/list 與 tools/call。工具包含 initialize_environment、sample_point、learn_impacts、run_window_direct 和 rank_actions。",
                "initialize 負責註冊 scenario、baseline、外部邊界、設備、家具、時間與 estimator。冷氣設備狀態不是只有模式，也包含目標溫度、風速或 fan strength、水平與垂直出風角度，以及 fixed/swing 擺動設定。sample_point 查詢指定座標的溫濕照度估計；rank_actions 則在給定 sample 與三因子目標後，用包含這些 AC 操作參數的候選動作做反事實排序。",
                "所以本研究對 MCP 的定位是系統整合與工具化封裝：證明這個數位孿生模型可以被 AI client 操作。我的研究貢獻不是提出新的 MCP protocol，也不是宣稱模型權重原生支援 MCP。",
                "Web demo 負責人機互動展示，Gemma/Ollama bridge 負責把自然語言轉成 tool calling。兩者底層都呼叫同一個模型服務，因此結果可以保持一致。",
            ],
        ),
        (
            "learn_impacts：動作如何成為資料記錄",
            [
                "這一頁回答「學習時記錄的動作是什麼」。系統不是只記錄一個開關，也不是記錄 rank_actions 的推薦名稱，而是把實際要套用到裝置上的 device_state 記錄下來。",
                "在 start 階段，client 會送 device_name 與 device_state。以冷氣為例，device_state 可以包含 activation、kind、power、ac_mode、target_temperature、fan_speed、fan_strength、horizontal_mode、horizontal_angle_deg、vertical_mode、vertical_angle_deg，以及 swing 相關欄位。系統會把這些欄位合併到目前註冊設備，形成新的 device_specs，並更新 runtime state。",
                "同一筆 learning record 會得到 learning_record_id，狀態先是 RECORDING。紀錄裡會保存 device_name、device_state、合併後的 device_specs、當時的室內 baseline、外部邊界、家具遮蔽、elapsed time、sampling mode、before_observations，以及 optional note。before_observations 的格式是 sensor name 對應 temperature、humidity、illuminance，例如 floor_sw 對應一組 T/H/L。",
                "finish 階段必須提供同一批感測器的 after_observations。系統用 after minus before 得到每顆感測器的 ΔT、ΔH、ΔL，再用模型中的 influence envelope 當作 X 矩陣，對每個 metric 解 least-squares 係數。輸出 learned_device_impacts 裡會包含 metric_coefficients、sensor_mae 與 sensor_observation_delta。",
                "例如一次冷氣學習事件可以記錄為：冷氣模式 cool、target_temperature 24°C、fan_strength high、horizontal_angle 30°、vertical_angle -15°、swing disabled。系統先保存操作前 8 顆感測器的 T/H/L，再等待操作後同一批感測器的 T/H/L。",
                "整理來說，學習資料是「操作狀態 + 環境快照 + 前後感測讀值」形成的事件紀錄；只有 before 和 after 都存在時才會計算係數。若想追蹤這筆資料來自哪一個推薦動作，目前可寫在 note，或未來新增 action_name 欄位。",
            ],
        ),
        (
            "驗證設計",
            [
                "驗證採分層設計，因為不同資料能支持的結論不同。E1 到 E3 使用受控完整場資料，重點是直接比較整個 3D field 的估計誤差，並和 IDW baseline 及 ablation variants 比較。",
                f"E4 到 E6 驗證模型的其他元件。E4 檢查非連網裝置影響學習是否能從 before/after observations 解出合理係數；E5 檢查 {window_summary.get('count', 0)} 組窗戶矩陣與 direct input 對外部邊界的支援，其中 {window_in_domain} 組 target-zone 室內溫度位於 20–30 °C，{window_out_of_domain} 組只作範圍外壓力測試；E6 檢查 hybrid no-Fourier 和 leave-one-scenario-out，確認改善不是單一切分造成。",
                f"E7 使用 bedroom_01 的 {bedroom_summary['snapshot_count']} 筆真實快照做 pillow hold-out 檢查。這裡能支持的是 sparse real-room calibration，也就是校正後對未參與 fitting 的 pillow 點有改善，但它不是 dense 3D truth。",
                "E8 是推薦動作的 before/after intervention protocol。也就是說，現在系統可以做 counterfactual ranking，但實際採取推薦後是否真的降低 comfort penalty，仍需要用介入實驗補上因果驗證。",
                "E9 使用公開資料集做 task-aligned benchmark。這部分不是單房間完整場驗證，而是把公開資料中相容的 boundary-response 或 device-response 子任務拿來壓力測試模型概念。",
            ],
        ),
        (
            "證據鏈與驗證範圍",
            [
                "這頁是把證據鏈和結論邊界講清楚。第一類證據是 synthetic full-field，因為它有完整 3D truth，所以可以真正計算整個房間的 field MAE，並比較 base model、IDW、ablation 和 hybrid residual。",
                "第二類證據是 real-bedroom snapshot。它來自真實臥室，因此比 synthetic 更貼近實際環境；但它只有稀疏量測與 pillow hold-out，不是整個房間每個位置都有 truth。因此它支援的是校正效果檢查，不是完整真實 3D 場驗證。",
                "第三類證據是 public datasets。SML2010 和 CU-BEMS 可以測某些相容子任務，例如外部邊界響應或設備響應，但它們沒有本研究的房間幾何、8 角點感測拓樸與 dense field truth。",
                "第四類是 recommendation。現階段推薦排序是 counterfactual simulation，也就是模型假設採取某動作後重新估計結果；它可以說明系統具備排序能力，但還不能宣稱推薦已被真實介入證明有效。",
                "因此我的結論會分層表述：受控完整場支援主要模型有效性，真實臥室支援 sparse calibration 的可行性，公開資料支援相容任務上的外部壓力測試，推薦則是已完成 protocol 與系統流程、仍待實測介入驗證。",
            ],
        ),
        (
            "情境設計與輸入模式",
            [
                "標準情境共有 8 組，包含 idle、ac_only、window_only、light_only、ac_window、window_light、ac_light 和 all_active。這些情境讓模型可以分別觀察單一設備、兩兩組合，以及全部設備同時作用時的三因子場變化。",
                "這種設計的目的不是列出所有可能生活情境，而是建立可控的 benchmark family。單裝置情境可以看每個裝置的基本影響，雙裝置情境可以看交互作用，all_active 則檢查多來源影響疊加時模型是否仍穩定。",
                f"窗戶相關輸入除了標準情境，也包含 {window_summary.get('count', 0)} 組 window matrix。矩陣會組合季節、天氣、時段等條件；依 target-zone 室內溫度稽核，{window_in_domain} 組位於 20–30 °C，另 {window_out_of_domain} 組只保留為範圍外壓力測試。",
                "系統也支援 direct input，讓使用者不一定要選預設矩陣，而可以直接輸入外部溫度、濕度、日照與開窗比例。這對 demo 和 MCP tools 很重要，因為使用者可能會問一個當下的自訂條件。",
                "所有 scenario 都有 elapsed time，設備影響用一階收斂近似描述從剛啟動到接近 quasi-steady state 的過程。這讓系統不只看靜態開關，也能呈現時間推進後的環境變化。",
            ],
        ),
        (
            "主要量化結果",
            [
                f"這頁上方的柱狀圖是 field_mae_comparison，資料來自 {scenario_count} 組 canonical scenarios 的完整 3D grid 評估。每一個 scenario 都會在整個房間網格上比較估計值與 truth，再把 temperature、humidity、illuminance 各自的 field MAE 平均起來。",
                "圖表的三個群組分別代表溫度、濕度與照度。Y 軸是 log-scale，原因是照度 MAE 的量級比溫度與濕度大很多；因此這張圖要看柱上數字與相對排序，不能只用柱高做線性比例解讀。",
                "每個群組內有三根柱。IDW 使用同一批 8 個角落觀測值做反距離插值，只知道距離與感測值；Base 是本研究的可解釋主模型，包含變數專屬 nominal model、裝置與幾何先驗、power calibration 與 trilinear residual correction；LOO Hybrid 是 leave-one-scenario-out 的 residual model 平均結果，用來檢查第二層 residual 是否只對單一切分有效。",
                f"具體數字上，IDW 的平均 field MAE 是 {metric_triplet(idw_mae)}；Base 是 {metric_triplet(base_mae)}；LOO Hybrid 是 {metric_triplet(loo_hybrid_mae)}。Base 相對 IDW 的降幅約為 temperature {percent_reduction(idw_mae['temperature'], base_mae['temperature']):.1f}%、humidity {percent_reduction(idw_mae['humidity'], base_mae['humidity']):.1f}%、illuminance {percent_reduction(idw_mae['illuminance'], base_mae['illuminance']):.1f}%。",
                "解讀重點是：IDW 在照度特別差，因為它不知道窗戶日照方向、燈具位置、家具遮蔽或反射；只靠距離很難重建局部光照分布。溫度與濕度也有改善，表示設備狀態、方向性與房間幾何先驗確實提供了純幾何插值沒有的資訊。",
                f"右下角的真實臥室校正檢查不是同一張柱狀圖的資料，而是 E7 bedroom_01 的 {bedroom_summary['snapshot_count']} 筆真實快照，房間網格解析度為 {grid_resolution['nx']} x {grid_resolution['ny']} x {grid_resolution['nz']}。Pillow 參考點沒有參與 8 角點 residual fitting，所以可當 held-out point 檢查非感測點估計。",
                f"真實臥室 pillow 點的 raw MAE 為 {metric_triplet(bedroom_aggregate['raw_pillow_mae'])}，校正後 MAE 為 {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}。相對 raw，校正後降幅約為 temperature {percent_reduction(bedroom_aggregate['raw_pillow_mae']['temperature'], bedroom_aggregate['estimated_pillow_mae']['temperature']):.1f}%、humidity {percent_reduction(bedroom_aggregate['raw_pillow_mae']['humidity'], bedroom_aggregate['estimated_pillow_mae']['humidity']):.1f}%、illuminance {percent_reduction(bedroom_aggregate['raw_pillow_mae']['illuminance'], bedroom_aggregate['estimated_pillow_mae']['illuminance']):.1f}%。這支持稀疏校正在此真實快照設定下有改善，但仍不能宣稱已具備 dense real-room ground truth 驗證。",
                f"另外，E7 不是把 28 筆快照全部當成互相獨立，而是用日期作為 block，把同一天四個時段一起重抽樣。固定 seed 執行 {bedroom_bootstrap['replicates']:,} 次後，temperature、humidity、illuminance 的 MAE 降幅 95% CI 分別為 [{bedroom_bootstrap_metrics['temperature']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['temperature']['ci95_absolute_mae_reduction']['upper']:.4f}]、[{bedroom_bootstrap_metrics['humidity']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['humidity']['ci95_absolute_mae_reduction']['upper']:.4f}]、[{bedroom_bootstrap_metrics['illuminance']['ci95_absolute_mae_reduction']['lower']:.4f}, {bedroom_bootstrap_metrics['illuminance']['ci95_absolute_mae_reduction']['upper']:.4f}]，下界都大於零。",
                f"再逐一移除七個日期中的任一天後，三因子最小 MAE 降幅仍為 {bedroom_lodo_metrics['temperature']['minimum_absolute_mae_reduction']:.4f}、{bedroom_lodo_metrics['humidity']['minimum_absolute_mae_reduction']:.4f} 與 {bedroom_lodo_metrics['illuminance']['minimum_absolute_mae_reduction']:.4f}。這表示正向改善不依賴保留某一特定日期，但 folds 高度重疊，仍不是獨立重複實驗。",
            ],
        ),
        (
            "真實臥室快照與推薦驗證狀態",
            [
                "E7 的重點是 pillow 參考點沒有參與 8 個角點 residual fitting，因此它可以用來檢查校正場是否改善非感測點估計。",
                f"結果上，校正後 pillow MAE 從 raw 的 {metric_triplet(bedroom_aggregate['raw_pillow_mae'])} 降到 {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}。",
                f"為保留同一天四個時段的相依性，我用 7 個日期作 block 做 {bedroom_bootstrap['replicates']:,} 次 paired bootstrap。逐快照改善數是 temperature {bedroom_bootstrap_metrics['temperature']['snapshots_improved']}/28、humidity 28/28、illuminance 28/28；這只能解讀成七天內的穩定改善，不能叫做控制介入成功率。",
                f"逐日剔除分析進一步顯示，移除任一天後的最小 MAE 降幅仍為 T {bedroom_lodo_metrics['temperature']['minimum_absolute_mae_reduction']:.4f}、H {bedroom_lodo_metrics['humidity']['minimum_absolute_mae_reduction']:.4f}、L {bedroom_lodo_metrics['illuminance']['minimum_absolute_mae_reduction']:.4f}，所以結果不依賴保留某一特定日期。",
                "這個 bootstrap 沒有增加新的獨立資料，所以 E7 的外部效度限制不變：仍是單一房間、單一 pillow hold-out，而且沒有 dense 3D ground truth。",
                "E8 現在不只是一段文字 protocol：repository 已有 versioned JSON schema、bedroom_01 pillow 的空白 trial template，以及會重新計算 penalty、actual improvement、prediction error、direction accuracy、top-1 regret 與 rank correlation 的 deterministic analyzer。",
                f"目前 machine-readable summary 顯示 completed real intervention trials 為 {e8_summary['trial_counts']['completed']}，status 是 {e8_summary['evidence_status']}，所有 efficacy estimates 都是 null。這代表可以直接開始收資料，但不能把 synthetic unit tests 或空白模板說成推薦效果。",
            ],
        ),
        (
            "3D 視覺化結果",
            [
                "這頁展示不同情境下的 3D 場分布，例如 all_active 的溫度場、window_only 的照度場，以及 ac_only 的溫度場。",
                "3D 圖的功能是幫助理解空間分布，不直接作為新的量化實驗。量化結果仍以前面提到的 field MAE、baseline comparison 和 hold-out 檢查為主。",
                "可以看到裝置位置與作用方向會造成局部差異，這也是為什麼模型不能只用單一平均值或純距離插值處理。",
            ],
        ),
        (
            "Hybrid Residual 結果",
            [
                f"Default held-out 的 train/test samples 為 {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}，hybrid test MAE 為 {metric_triplet(default_hybrid['hybrid_test_field_mae'])}。",
                f"No-Fourier 對照的 MAE 為 {metric_triplet(no_fourier['hybrid_test_field_mae'])}，LOO 平均 hybrid MAE 為 {metric_triplet(loo['average_hybrid_field_mae'])}。",
                "這些結果表示標準情境 family 內的 residual 有可學習性，但不能直接擴大解讀為任意房間、任意裝置配置都能泛化。",
            ],
        ),
        (
            "公開資料任務拆解：SML2010",
            [
                "SML2010 被映射成 two-point boundary-response benchmark。它適合檢查外氣、日照與室內兩點響應，但沒有完整 3D 場真值。",
                "S1 純照度短視窗是主要劣勢，因為 persistence 在短時間照度高度自相關時很強。S2 長視窗溫度有部分優勢，但濕度有尺度對齊問題。",
                "S3 facade event delta 是主要優勢，因為事件後變化方向和長視窗響應更能受益於 structured prior。",
                "另外把 Oh et al. 2024 的 simulation-plus-residual 概念移植成固定 ridge-linear residual head。15 分鐘兩個溫度點由 transfer 最佳，60 分鐘由本研究 readout 最佳，24 小時則由 persistence 最佳；transfer 在 24 小時還劣於 raw physics。",
                "為嘗試增加次日優勢，primary follow-up 用 60/10/30 chronological split，只在 validation 選模型。兩點都選到 alpha=0.25 daily trend，但 final test 反而比 persistence 惡化 7.34% 與 8.36%，date-block bootstrap interval 也都跨 0。",
                "Primary 結果後另做明確標記的 adaptive online exploratory analysis，validation 選到 14-day median，但 test 仍惡化 8.83% 與 9.73%。Registered bias correction 雖約改善 1%，卻未被 validation 選中，因此不能事後包裝成次日優勢。",
                "依教授建議，我也加入固定 vanilla RNN。四種方法先共用完全相同的四筆歷史、chronological split、targets 與 test rows，並以 endpoint 和 input-content hash 稽核。12 個案例全部通過 parity；最低 MAE 由 sequence linear regression 取得 7 項、persistence 取得 5 項，RNN 為 0 項。",
                "這個結果必須保留。它代表在目前資料、四筆歷史與固定小型架構下，recurrent complexity 沒有帶來可驗證優勢；不能因為 RNN 較複雜就把它當成改進後模型。",
                "原文 BEMS data 為 confidential，且沒有可用的 CNN--LSTM code，因此這只能稱 published-method-inspired transfer，不能稱重現原文 next-day performance。",
            ],
        ),
        (
            "公開資料任務拆解：CU-BEMS",
            [
                "CU-BEMS 被映射成商辦 zone-level device-response benchmark。它有 AC power 和 lighting power 等欄位，但不是本研究的單房間 8 點拓樸。",
                "C1 中 AC 溫濕度可補強 linear regression，但不勝 persistence。C2 商辦照度與單房間光學假設差距大，是明確劣勢。",
                "C3 compound event 可勝 linear regression，但仍不勝 persistence。這表示本研究特徵對事件讀出有幫助，但不能宣稱在商辦時序任務全面勝出。",
            ],
        ),
        (
            "結論、限制與未來工作",
            [
                "本研究完成一個單房間三因子空間數位孿生原型，能在少量角落感測器下估計溫度、濕度與照度分布，並支援非連網裝置影響學習。",
                "限制方面，目前仍缺長期 dense real-room ground truth，hybrid residual 的泛化也主要限於標準情境 family。",
                "教授提醒後，室內溫度適用範圍明確限制在 20–30 °C，人體舒適改以目標帶和容許範圍判定；現有低 MAE 不能直接證明一般人居空間需要極窄控制，也不能外推到超出溫度範圍的用途。",
                "需要動態環境配方的小型封閉植物生長空間可作候選，但目前 lux 不是 PPFD/PAR，並缺 CO2、基質水分、氣流與生物 endpoint，所以不能宣稱已具培養成效。",
                "Kalman filter 目前只列為狀態估測、感測融合或線上參數調整的後續方法。未來要先定義 state、observation 與 covariance，再讓未濾波、moving average、KF 與 EKF 使用相同資料比較。",
                "其他未來工作包括擴大 ESP32 長期資料、發展 multi-zone model、執行推薦動作介入驗證，以及往閉環控制延伸。",
            ],
        ),
        (
            "公式與指標整理",
            [
                "後半段整理公式與指標。第一組是場模型，包括三因子場、總估計式、baseline、activation 與 influence envelope。",
                "第二組是三因子 nominal model，分別說明溫度、濕度與照度為什麼要採用不同的物理近似。",
                "第三組是校正與評估，包括 8 點三線性 residual correction、非連網裝置影響學習、hybrid residual、MAE/RMSE/correlation、IDW baseline 與推薦排序。",
            ],
        ),
    ]

    for title, left_title, left_lines, right_title, right_lines in FORMULA_WALKTHROUGH:
        example = FORMULA_NUMERIC_EXAMPLES.get(title)
        if title == "公式說明 22：平滑 residual 的誤差界":
            paragraphs = [
                "這頁可以用很白話的方式說：我們只有 8 個角落點有真實 residual，房間中間很多點沒有直接量到，所以 Cᵥ 是用 8 個角點補出來的 residual。公式左邊 |Rᵥ-Cᵥ| 就是在問：某個沒量到的點，這個補出來的 residual 最多可能和真實 residual 差多少。",
                "M_xx、M_yy、M_zz 的名字來自二階偏導數記號。M_xx 不是 x 乘 x，而是對 x 方向微分兩次後的最大絕對值上界，也就是 M_xx ≥ max|∂²Rᵥ/∂x²|。M_yy 和 M_zz 同理，分別代表 y 與 z 方向 residual 的最大彎曲程度。",
                "右邊分成三個方向：W²M_xx/8 是 x 方向造成的最壞誤差，L²M_yy/8 是 y 方向，H²M_zz/8 是 z 方向。W、L、H 越大，代表角點之間隔得越遠，中間靠補間猜的距離越長；M_xx、M_yy、M_zz 越大，代表 residual 在該方向彎得越厲害，也越難用直線補準。",
                "原因可以用「用直線補曲線」來理解。線性補間等於拿兩端點連成一條直線去估中間值；如果 residual 是直線，即使斜率很大，二階導數仍然是 0，線性補間可以補對；真正造成補間誤差的是曲線彎曲，也就是二階導數。這就是為什麼這裡用 M_xx、M_yy、M_zz，而不是用一階斜率。",
                "三線性補間只是把這件事放到 3D 房間裡：先沿 x 方向補，再沿 y 方向補，再沿 z 方向補。因此三個方向各自有一個可能誤差，合起來就是 W²M_xx/8 + L²M_yy/8 + H²M_zz/8。這也說明為什麼主模型要先把冷氣、窗戶、燈具等主要效果吃掉；剩下的 residual 越平滑，這個上界才越有意義。若 residual 是局部尖峰、光斑或遮蔽邊界，曲率會變大，單靠 8 點就不夠，需要更多感測點或 hybrid residual 補強。",
            ]
            if example:
                paragraphs.append(f"數字範例：{example[1]}")
            slides.append((title, paragraphs))
        else:
            left_text = "；".join(left_lines[:3])
            right_text = "；".join(right_lines[:3])
            paragraphs = [
                f"這頁說明「{left_title}」。可以先從公式或定義開始，指出它在整體模型中負責哪一部分。",
                f"左側重點包含：{left_text}。",
                f"接著說明「{right_title}」。這一部分通常用來補上模型設計理由、限制或可主張範圍。",
                f"右側重點包含：{right_text}。",
            ]
            if example:
                paragraphs.append(f"數字範例：{example[1]}")
            slides.append(
                (
                    title,
                    paragraphs,
                )
            )

    lines = [
        "# 30 分鐘論文簡報逐頁講稿",
        "",
        "本檔是 `thesis_presentation_zh_30min.pptx` 的講稿，不放入投影片畫面。投影片維持正式內容；這份 Markdown 用於練習口頭說明與答辯準備。",
        "",
    ]
    for index, (title, paragraphs) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        for paragraph in paragraphs:
            lines.append("")
            lines.append(paragraph)
        glossary_notes = glossary_notes_for_slide(title, paragraphs)
        if glossary_notes:
            lines.append("")
            lines.append("### 名詞註釋")
            for label, explanation in glossary_notes:
                lines.append(f"- **{label}**：{explanation}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def main() -> None:
    PAPERS.mkdir(parents=True, exist_ok=True)
    THESIS_PAPERS.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(PRESENTATION_PATH)
    prs.save(STORED_PRESENTATION_PATH)
    prs_long = build_presentation_30min()
    prs_long.save(LONG_PRESENTATION_PATH)
    prs_long.save(STORED_LONG_PRESENTATION_PATH)
    OUTLINE_PATH.write_text(build_outline(), encoding="utf-8")
    LONG_OUTLINE_PATH.write_text(build_outline_30min(), encoding="utf-8")
    LONG_SPEAKER_NOTES_PATH.write_text(build_speaker_notes_30min(), encoding="utf-8")
    print(f"Wrote {PRESENTATION_PATH}")
    print(f"Wrote {STORED_PRESENTATION_PATH}")
    print(f"Wrote {LONG_PRESENTATION_PATH}")
    print(f"Wrote {STORED_LONG_PRESENTATION_PATH}")
    print(f"Wrote {OUTLINE_PATH}")
    print(f"Wrote {LONG_OUTLINE_PATH}")
    print(f"Wrote {LONG_SPEAKER_NOTES_PATH}")


if __name__ == "__main__":
    main()
