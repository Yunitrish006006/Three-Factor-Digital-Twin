#!/usr/bin/env python3
"""Append synchronized E11G evidence to generated thesis and presentation files."""

from __future__ import annotations

import html
import tempfile
import zipfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
TOKEN = "E11G tail-safe"
DECISION = "no_candidate_forwarded"

DOCX_PATHS = (
    ROOT / "docs/papers/thesis/thesis_draft_zh.docx",
    ROOT / "outputs/papers/thesis_draft_zh.docx",
)
THESIS_MARKDOWN = ROOT / "docs/thesis/thesis_draft_zh.md"
PPTX_PATHS = (
    ROOT / "outputs/papers/thesis_presentation_zh.pptx",
    ROOT / "outputs/papers/thesis_presentation_zh_30min.pptx",
)
OUTLINE_PATHS = (
    ROOT / "docs/thesis/presentation_outline_zh.md",
    ROOT / "docs/thesis/presentation_outline_zh_30min.md",
)


def _docx_paragraph(text: str, bold: bool = False) -> str:
    run_properties = "<w:rPr><w:b/></w:rPr>" if bold else ""
    return (
        "<w:p><w:r>"
        + run_properties
        + '<w:t xml:space="preserve">'
        + html.escape(text)
        + "</w:t></w:r></w:p>"
    )


def _append_docx(path: Path) -> None:
    if not path.exists():
        return
    with zipfile.ZipFile(path, "r") as archive:
        entries = {item.filename: archive.read(item.filename) for item in archive.infolist()}
    document = entries["word/document.xml"].decode("utf-8")
    if TOKEN in document:
        return
    addition = "".join(
        (
            _docx_paragraph("E11G tail-safe 自適應開發結果", bold=True),
            _docx_paragraph(
                "E11G 在 E11E 開發資料上進行 12 日 leave-one-day-out 與折內感測器選擇。"
                "相較 local-IDW，MAE 由 1.1168°C 降至 0.8945°C，RMSE 由 1.7250°C 降至 1.5415°C，"
                "P95 由 3.4900°C 降至 3.1013°C；日區塊 MAE 改善 95% CI 為 [0.1847, 0.2620]°C。"
            ),
            _docx_paragraph(
                "但嚴格感測器勝率只有 21/42，未達預註冊 26/42；20 個感測器因安全回退而持平，"
                "1 個微幅惡化。因此決策為 no_candidate_forwarded，E11F 未存取，結果僅屬適應性開發證據。"
            ),
        )
    )
    entries["word/document.xml"] = document.replace("</w:body>", addition + "</w:body>").encode("utf-8")
    with tempfile.NamedTemporaryFile(dir=path.parent, suffix=".docx", delete=False) as temporary:
        temporary_path = Path(temporary.name)
    try:
        with zipfile.ZipFile(temporary_path, "w", zipfile.ZIP_DEFLATED) as archive:
            for name, content in entries.items():
                archive.writestr(name, content)
        temporary_path.replace(path)
    finally:
        temporary_path.unlink(missing_ok=True)


def _append_thesis_markdown() -> None:
    if not THESIS_MARKDOWN.exists():
        return
    content = THESIS_MARKDOWN.read_text(encoding="utf-8")
    if TOKEN in content:
        return
    content += (
        "\n\n## E11G tail-safe 自適應開發與空間覆蓋限制\n\n"
        "E11G 針對 E11E 的 P95 惡化問題，在同一開發資料上預註冊 12 日 leave-one-day-out。"
        "方法以 local-IDW（k=3、p=2）為安全基線，以同角色 local-IDW（k=5、p=2）為專家，"
        "評估 30 組裁切修正與高分歧回退規格。每一折的感測器選擇只能使用其餘日期，"
        "候選必須同時改善 MAE、RMSE、P95 至少 0.02°C，並在至少 60% 訓練日降低 MAE；"
        "否則測試日回退基線。\n\n"
        "在 42 個感測器、12 日與 63,084 筆感測器分鐘評估上，out-of-fold MAE 由 1.1168°C "
        "降至 0.8945°C，RMSE 由 1.7250°C 降至 1.5415°C，P95 由 3.4900°C 降至 3.1013°C；"
        "日區塊 MAE 改善 bootstrap 95% CI 為 [0.1847, 0.2620]°C。然而嚴格感測器勝率僅 "
        "21/42，另有 20 個因安全回退而持平、1 個微幅惡化，未達預註冊 26/42。"
        "因此決策為 `no_candidate_forwarded`，E11F 未存取；本結果僅屬適應性開發證據。\n"
    )
    THESIS_MARKDOWN.write_text(content, encoding="utf-8")


def sync_docx_outputs() -> None:
    _append_thesis_markdown()
    for path in DOCX_PATHS:
        _append_docx(path)


def _append_outline(path: Path) -> None:
    if not path.exists():
        return
    content = path.read_text(encoding="utf-8")
    if TOKEN in content:
        return
    content += (
        "\n\n## E11G tail-safe 自適應開發\n\n"
        "- 12 日 leave-one-day-out、42 感測器、30 個裁切與回退候選。\n"
        "- MAE 1.1168→0.8945°C；RMSE 1.7250→1.5415°C；P95 3.4900→3.1013°C。\n"
        "- bootstrap 95% CI：[0.1847, 0.2620]°C。\n"
        "- 嚴格勝率僅 21/42，低於 26/42；`no_candidate_forwarded`，E11F 未存取。\n"
    )
    path.write_text(content, encoding="utf-8")


def _append_pptx(path: Path) -> None:
    if not path.exists():
        return
    from pptx import Presentation

    presentation = Presentation(path)
    existing = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    if TOKEN in existing:
        return
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "E11G tail-safe：尾端改善但覆蓋未過"
    slide.placeholders[1].text = (
        "12 日 leave-one-day-out，42 感測器，30 候選\n"
        "MAE：1.1168 → 0.8945°C\n"
        "RMSE：1.7250 → 1.5415°C\n"
        "P95：3.4900 → 3.1013°C\n"
        "日區塊改善 95% CI：[0.1847, 0.2620]°C\n"
        "嚴格勝率 21/42 < 26/42；no_candidate_forwarded，E11F 未存取"
    )
    presentation.save(path)


def sync_pptx_outputs() -> None:
    for path in OUTLINE_PATHS:
        _append_outline(path)
    for path in PPTX_PATHS:
        _append_pptx(path)
