#!/usr/bin/env python3
"""Build the 2026-08-24 Chinese research progress presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs/reports/research_progress_2026-08-24_zh.pptx"

INK = RGBColor(28, 36, 38)
PAPER = RGBColor(245, 241, 232)
TEAL = RGBColor(16, 112, 108)
ORANGE = RGBColor(222, 106, 55)
MINT = RGBColor(207, 226, 216)
SAND = RGBColor(229, 216, 190)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(176, 55, 45)
GRAY = RGBColor(91, 101, 101)
FONT = "Noto Sans CJK TC"
DISPLAY = "Noto Serif CJK TC"


def add_text(slide, text, x, y, w, h, size=22, color=INK, bold=False,
             font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.level = 0
        paragraph.space_after = Pt(12)
        paragraph.text = "•  " + paragraph.text
    return box


def add_rect(slide, x, y, w, h, fill, radius=True, line=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def base_slide(prs, number, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PAPER
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL, radius=False)
    add_text(slide, section.upper(), 0.55, 0.22, 5.0, 0.3, 10, TEAL, True)
    add_text(slide, f"{number:02d}", 12.25, 0.18, 0.5, 0.35, 12, ORANGE, True,
             align=PP_ALIGN.RIGHT)
    add_text(slide, "2026.08.24  |  機箱溫度虛擬感測", 0.55, 7.08, 5.5, 0.22, 9, GRAY)
    return slide


def title(slide, heading, sub=None):
    add_text(slide, heading, 0.72, 0.72, 11.8, 0.72, 30, INK, True, DISPLAY)
    if sub:
        add_text(slide, sub, 0.75, 1.47, 11.2, 0.44, 15, GRAY)


def card(slide, heading, body, x, y, w, h, fill=MINT, accent=TEAL):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.08, h, accent, radius=False)
    add_text(slide, heading, x + 0.25, y + 0.2, w - 0.45, 0.35, 16, accent, True)
    add_text(slide, body, x + 0.25, y + 0.68, w - 0.45, h - 0.85, 15, INK)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = base_slide(prs, 1, "Research progress")
    add_rect(slide, 8.7, 0.0, 4.63, 7.5, TEAL, radius=False)
    add_rect(slide, 9.45, 0.75, 2.9, 2.9, ORANGE)
    add_text(slide, "DATA\nBEFORE\nMODEL", 9.72, 1.16, 2.35, 1.95, 30, WHITE, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "機箱溫度\n虛擬感測研究進度", 0.85, 1.42, 7.25, 1.75, 38, INK, True, DISPLAY)
    add_text(slide, "從失敗的跨 run 實驗，到可確認的凍結候選模型", 0.88, 3.48, 7.0, 0.55, 20, TEAL, True)
    add_text(slide, "今日重點：資料完整性、單位制度、E14C 結果、E15 設計", 0.88, 4.3, 7.1, 0.75, 17, GRAY)
    add_text(slide, "報告者：Thomas", 0.88, 5.7, 4.0, 0.4, 15, INK, True)

    slide = base_slide(prs, 2, "Question")
    title(slide, "現在已證明能用在電腦機箱嗎？")
    add_rect(slide, 0.78, 2.05, 3.2, 3.75, SAND)
    add_text(slide, "還沒有", 1.12, 2.6, 2.5, 0.7, 37, RED, True, DISPLAY,
             align=PP_ALIGN.CENTER)
    add_text(slide, "目前是公開伺服器 BMC 的同來源證據，不是桌機、NTC 或跨硬體證據。",
             1.1, 3.65, 2.55, 1.35, 17, INK, align=PP_ALIGN.CENTER)
    card(slide, "今天真正完成的事", "把不可相信的跨 run 結果拆解，修正 source、header 與 unit 問題，建立能接受新資料檢驗的候選模型。",
         4.45, 2.05, 7.65, 1.55, MINT, TEAL)
    card(slide, "研究定位", "稀疏 IoT 感測 → 虛擬溫度感測 → 空間智慧與決策支援。MCP 只作為後續服務層。",
         4.45, 4.08, 7.65, 1.55, WHITE, ORANGE)

    slide = base_slide(prs, 3, "Evidence chain")
    title(slide, "證據不是一次得到，而是逐關卡建立")
    stages = [
        ("E11H/F", "同 campaign\n初步正結果", TEAL),
        ("E12", "30-row 閘門\n提前失敗", ORANGE),
        ("E13", "極端數值\nparser-invalidated", RED),
        ("E14A/B", "來源與單位\n資料管線修正", TEAL),
        ("E14C", "回溯敏感度\n候選模型通過", ORANGE),
        ("E15", "新 runs\n等待確認", INK),
    ]
    for i, (label, body, color) in enumerate(stages):
        x = 0.72 + i * 2.05
        add_rect(slide, x, 2.35, 1.72, 2.22, WHITE, line=color)
        add_rect(slide, x, 2.35, 1.72, 0.48, color, radius=False)
        add_text(slide, label, x + 0.1, 2.42, 1.52, 0.25, 14, WHITE, True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.15, 3.12, 1.42, 0.95, 15, INK, True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            add_text(slide, "→", x + 1.75, 3.13, 0.3, 0.45, 22, GRAY, True)
    add_text(slide, "原則：任何資料修正都不能把已開啟的測試重新宣稱為 unseen。",
             1.02, 5.35, 11.2, 0.55, 19, RED, True, align=PP_ALIGN.CENTER)

    slide = base_slide(prs, 4, "Data integrity")
    title(slide, "真正的難題先發生在資料，不在模型")
    card(slide, "1  Header 漂移", "每個 InfluxDB #group section 可能有不同欄位。固定第一個 header 會把 host counter 誤當 BMC temperature。",
         0.75, 2.0, 3.7, 2.05, WHITE, RED)
    card(slide, "2  Source 混合", "同檔案同時含 BMC 與 host。E14A 只接受 measurement=sdgp、device_id=bmc。",
         4.82, 2.0, 3.7, 2.05, WHITE, ORANGE)
    card(slide, "3  Unit 混合", "三檔使用 millidegree C 與 microwatt。E14B 逐 section 判定縮放制度。",
         8.89, 2.0, 3.7, 2.05, WHITE, TEAL)
    add_rect(slide, 1.28, 4.72, 10.75, 1.15, INK)
    add_text(slide, "E14B 結果：31/31 檔案一致  |  4,038 rows  |  8/8 資料品質閘門通過",
             1.62, 5.03, 10.05, 0.42, 20, WHITE, True, align=PP_ALIGN.CENTER)

    slide = base_slide(prs, 5, "E14C result")
    title(slide, "修正資料後，低複雜度模型仍明顯改善")
    metrics = [("MAE", 4.088, 1.805), ("RMSE", 5.209, 2.800), ("P95", 12.000, 7.115)]
    max_value = 12.0
    for index, (name, baseline, ridge) in enumerate(metrics):
        y = 2.05 + index * 1.32
        add_text(slide, name, 0.82, y + 0.22, 0.75, 0.35, 16, INK, True)
        add_rect(slide, 1.72, y, 7.7 * baseline / max_value, 0.38, SAND, radius=False)
        add_rect(slide, 1.72, y + 0.5, 7.7 * ridge / max_value, 0.38, TEAL, radius=False)
        add_text(slide, f"{baseline:.3f}", 9.55, y - 0.02, 0.85, 0.32, 14, INK, True)
        add_text(slide, f"{ridge:.3f}", 9.55, y + 0.48, 0.85, 0.32, 14, TEAL, True)
    add_text(slide, "Baseline", 10.75, 2.08, 1.0, 0.28, 13, INK, True)
    add_rect(slide, 10.45, 2.1, 0.2, 0.2, SAND, radius=False)
    add_text(slide, "Ridge", 10.75, 2.58, 1.0, 0.28, 13, TEAL, True)
    add_rect(slide, 10.45, 2.6, 0.2, 0.2, TEAL, radius=False)
    card(slide, "13 / 14", "runs 勝過 baseline", 10.35, 3.42, 2.15, 1.12, MINT, TEAL)
    card(slide, "[1.427, 2.794] °C", "macro MAE gain 95% CI", 10.35, 4.8, 2.15, 1.12, SAND, ORANGE)
    add_text(slide, "14 runs  |  3,242 rows  |  retrospective sensitivity, not independent confirmation",
             1.75, 6.2, 9.6, 0.4, 15, RED, True, align=PP_ALIGN.CENTER)

    slide = base_slide(prs, 6, "Claim boundary")
    title(slide, "研究價值來自清楚的證據邊界")
    add_rect(slide, 0.78, 1.9, 5.75, 4.5, MINT)
    add_text(slide, "現在可以說", 1.12, 2.18, 4.9, 0.45, 22, TEAL, True, DISPLAY)
    add_bullets(slide, [
        "資料管線已能辨識 section、source 與 unit regime",
        "load-aware ridge 是值得確認的候選模型",
        "E14C 的平均、尾端與跨 run 指標都優於 baseline",
    ], 1.1, 2.9, 4.95, 2.75, 17)
    add_rect(slide, 6.82, 1.9, 5.75, 4.5, SAND)
    add_text(slide, "現在不能說", 7.16, 2.18, 4.9, 0.45, 22, RED, True, DISPLAY)
    add_bullets(slide, [
        "不能宣稱已驗證桌上型電腦機箱",
        "不能宣稱跨伺服器、跨感測器泛化",
        "沒有 NTC 實體真值或空間溫度場證據",
        "E14C 不是獨立 unseen confirmation",
    ], 7.14, 2.9, 4.95, 2.9, 17)

    slide = base_slide(prs, 7, "Hardware")
    title(slide, "NTC 可以用，但它不是插上就是真值")
    add_text(slide, "可行", 0.88, 1.95, 1.2, 0.55, 27, TEAL, True, DISPLAY)
    add_bullets(slide, [
        "低成本、體積小，適合多點機箱測溫",
        "分壓 + ADC + Beta / Steinhart–Hart 可轉換溫度",
        "適合作為後續稀疏 IoT 感測節點",
    ], 0.9, 2.6, 5.25, 2.6, 18)
    add_text(slide, "必要條件", 6.72, 1.95, 2.0, 0.55, 27, ORANGE, True, DISPLAY)
    add_bullets(slide, [
        "參考溫度計與多點校正",
        "精密電阻、穩定 ADC 與參考電壓",
        "控制自熱、導線電阻、安裝接觸與熱慣性",
        "固定感測點座標與量測 protocol",
    ], 6.74, 2.6, 5.55, 2.95, 18)
    add_rect(slide, 1.4, 5.77, 10.4, 0.72, INK)
    add_text(slide, "本輪決策：實體先不動；先完成 E15 資料驗證。",
             1.75, 5.94, 9.7, 0.32, 18, WHITE, True, align=PP_ALIGN.CENTER)

    slide = base_slide(prs, 8, "Research decisions")
    title(slide, "本週真正的大問題：研究主張與證據怎麼對齊")
    items = [
        ("研究對象", "BMC CPU hotspot ≠ 桌機機箱溫度場；先定位為方法階段證據"),
        ("泛化問題", "同 campaign 內插 ≠ 跨環境；改以完整 run 做時間與 workload 確認"),
        ("Target 定義", "區分元件熱點、機箱 zone、房間空間場，不用單一溫度取代全部"),
        ("模型策略", "獨立 runs 有限；先確認低維 ridge，再研究 robust / selective 方法"),
        ("論文主軸", "稀疏 IoT 感測與非聯網設備影響是 novelty；MCP 降為服務介面"),
    ]
    for index, (code, body) in enumerate(items):
        y = 1.9 + index * 0.92
        add_rect(slide, 0.82, y, 2.0, 0.62, TEAL if index % 2 == 0 else ORANGE)
        add_text(slide, code, 0.98, y + 0.13, 1.68, 0.28, 14, WHITE, True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, body, 3.15, y + 0.09, 9.0, 0.42, 16, INK, index == 2)
    add_text(slide, "技術 bug 與下載問題另列執行紀錄，不當作研究方向貢獻。",
             0.95, 6.55, 10.2, 0.3, 13, GRAY)

    slide = base_slide(prs, 9, "Next experiment")
    title(slide, "E15：真正的新資料凍結確認")
    card(slide, "固定模型", "E14C frozen model\nSHA-256: 609048…84", 0.8, 1.92, 3.65, 1.55, WHITE, TEAL)
    card(slide, "固定資料", "14 complete runs\n2023-08 → 2024-05", 4.84, 1.92, 3.65, 1.55, WHITE, ORANGE)
    card(slide, "固定判定", "10 conjunctive gates\n20,000 run-block bootstrap", 8.88, 1.92, 3.65, 1.55, WHITE, INK)
    add_bullets(slide, [
        "工作負載：For-And-Join、TEAP、kernel、synthetic、Jenkins",
        "禁止重訓、重新選特徵、改門檻或事後替換檔案",
        "通過也只支持同一伺服器內的跨時間與工作負載確認",
    ], 1.0, 4.05, 10.9, 1.8, 18)
    add_rect(slide, 3.25, 6.05, 6.8, 0.62, SAND)
    add_text(slide, "目前狀態：規格與測試通過，等待下載授權",
             3.48, 6.19, 6.35, 0.28, 16, RED, True, align=PP_ALIGN.CENTER)

    slide = base_slide(prs, 10, "Takeaway")
    add_text(slide, "今天的結論", 0.85, 1.05, 4.0, 0.58, 31, TEAL, True, DISPLAY)
    add_text(slide, "公開 BMC 是方法證據，不是最終機箱證據", 0.9, 2.0, 10.4, 0.62, 29, RED, True)
    add_text(slide, "研究路線必須逐步跨越 domain gap", 0.9, 2.88, 10.8, 0.75, 31, INK, True, DISPLAY)
    add_rect(slide, 0.9, 4.25, 11.45, 1.35, INK)
    add_text(slide, "BMC method → Temporal confirmation → PC / NTC zones → Spatial intelligence",
             1.18, 4.66, 10.9, 0.42, 20, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "下一步：取得 E15 資料，執行一次，無論成敗都保留。",
             1.1, 6.18, 10.8, 0.45, 19, ORANGE, True, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
