#!/usr/bin/env python3
"""Synchronize E11H development and E11F confirmation into generated artifacts."""

from __future__ import annotations

import html
import tempfile
import zipfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
TOKEN = "E11H commissioning"
DECISION = "h_enc_05_supported_within_campaign"
THESIS_MARKDOWN = ROOT / "docs/thesis/thesis_draft_zh.md"
DOCX_PATHS = (
    ROOT / "docs/papers/thesis/thesis_draft_zh.docx",
    ROOT / "outputs/papers/thesis_draft_zh.docx",
)
PPTX_PATHS = (
    ROOT / "outputs/papers/thesis_presentation_zh.pptx",
    ROOT / "outputs/papers/thesis_presentation_zh_30min.pptx",
)
OUTLINE_PATHS = (
    ROOT / "docs/thesis/presentation_outline_zh.md",
    ROOT / "docs/thesis/presentation_outline_zh_30min.md",
)


def _markdown_section() -> str:
    return (
        "\n\n## E11H commissioning 開發與 E11F 凍結確認\n\n"
        "E11H 將低成本 NTC 或參考感測器定位為短期 commissioning 工具：最早兩日的目標點真值"
        "用於 robust residual calibration，第三日只選擇模型，後續九日凍結測試。相較 local-IDW，"
        "E11H 測試 MAE 由 1.0958°C 降至 0.4039°C，RMSE 由 1.7435°C 降至 0.6830°C，P95 由 "
        "3.5061°C 降至 1.2900°C；39/42 感測器改善，日區塊 95% CI 為 [0.4854, 0.9271]°C。"
        "所有開發閘門通過，但多個 Huber slope 位於 0.5 邊界，窄溫域下 affine 參數不可解讀為物理係數。\n\n"
        "E11F 使用 11 個預先保留且未見的 byte ranges，完全凍結 E11H 的 42 組模型，不重新校正。"
        "相較 local-IDW，MAE 由 1.1399°C 降至 0.3966°C，RMSE 由 1.7850°C 降至 0.6723°C，"
        "P95 由 3.5735°C 降至 1.2756°C；39/42 感測器改善，13 日 bootstrap 95% CI 為 "
        "[0.5851, 0.9274]°C。因此 `h_enc_05_supported_within_campaign`。然而 11 個日期與 E11G "
        "重疊、8 個與 E11H 重疊，故證據只支持同一 AAU campaign 的 calibration-assisted unseen-byte "
        "transfer，不是跨日期、跨機箱、氣流因果或 NTC 硬體準確度驗證。\n"
    )


def _append_markdown() -> None:
    if not THESIS_MARKDOWN.exists():
        return
    content = THESIS_MARKDOWN.read_text(encoding="utf-8")
    if TOKEN not in content:
        THESIS_MARKDOWN.write_text(content + _markdown_section(), encoding="utf-8")


def _paragraph(text: str, bold: bool = False) -> str:
    properties = "<w:rPr><w:b/></w:rPr>" if bold else ""
    return (
        "<w:p><w:r>" + properties + '<w:t xml:space="preserve">'
        + html.escape(text) + "</w:t></w:r></w:p>"
    )


def _append_docx(path: Path) -> None:
    if not path.exists():
        return
    with zipfile.ZipFile(path) as archive:
        entries = {item.filename: archive.read(item.filename) for item in archive.infolist()}
    document = entries["word/document.xml"].decode("utf-8")
    if TOKEN in document:
        return
    addition = "".join(
        (
            _paragraph("E11H commissioning 開發與 E11F 凍結確認", bold=True),
            _paragraph(
                "E11H 採兩日校正、一日選模、九日凍結測試。MAE/RMSE/P95 由 "
                "1.0958/1.7435/3.5061°C 降至 0.4039/0.6830/1.2900°C，39/42 感測器改善，"
                "95% CI 為 [0.4854, 0.9271]°C；多個 Huber slope 達 0.5 邊界。"
            ),
            _paragraph(
                "E11F 不重新校正，MAE/RMSE/P95 為 0.3966/0.6723/1.2756°C，39/42 改善，"
                "95% CI 為 [0.5851, 0.9274]°C；h_enc_05_supported_within_campaign。"
                "日期與開發資料重疊，故不代表跨日期、跨機箱或 NTC 硬體驗證。"
            ),
        )
    )
    entries["word/document.xml"] = document.replace("</w:body>", addition + "</w:body>").encode("utf-8")
    with tempfile.NamedTemporaryFile(dir=path.parent, suffix=".docx", delete=False) as temporary:
        temporary_path = Path(temporary.name)
    try:
        with zipfile.ZipFile(temporary_path, "w", zipfile.ZIP_DEFLATED) as archive:
            for name, content in entries.items():
                archive.writestr(name, content)
        temporary_path.replace(path)
    finally:
        temporary_path.unlink(missing_ok=True)


def sync_docx_outputs() -> None:
    _append_markdown()
    for path in DOCX_PATHS:
        _append_docx(path)


def _append_outline(path: Path) -> None:
    if not path.exists():
        return
    content = path.read_text(encoding="utf-8")
    if TOKEN in content:
        return
    content += (
        "\n\n## E11H commissioning 與 E11F confirmation\n\n"
        "- E11H：2 日校正、1 日選模、9 日凍結測試；MAE 0.4039°C，P95 1.2900°C，39/42。\n"
        "- E11F：不 refit；MAE 0.3966°C、RMSE 0.6723°C、P95 1.2756°C，39/42。\n"
        "- H-ENC-05 僅於同 campaign 未見 bytes 獲支持；日期重疊，非跨機箱或 NTC 硬體驗證。\n"
        "- 決策：`h_enc_05_supported_within_campaign`。\n"
    )
    path.write_text(content, encoding="utf-8")


def _append_pptx(path: Path) -> None:
    if not path.exists():
        return
    from pptx import Presentation

    presentation = Presentation(path)
    existing = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    if TOKEN in existing:
        return
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "E11H commissioning：開發通過"
    slide.placeholders[1].text = (
        "2 日校正、1 日選模、9 日凍結測試\n"
        "MAE：1.0958 → 0.4039°C；P95：3.5061 → 1.2900°C\n"
        "39/42 感測器改善；95% CI：[0.4854, 0.9271]°C\n"
        "限制：Huber slope 邊界與 commissioning prerequisite"
    )
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "E11F：frozen no-refit confirmation"
    slide.placeholders[1].text = (
        "MAE/RMSE/P95：0.3966/0.6723/1.2756°C\n"
        "39/42 感測器改善；95% CI：[0.5851, 0.9274]°C\n"
        "h_enc_05_supported_within_campaign\n"
        "日期與開發重疊：僅同 campaign unseen-byte evidence\n"
        "不是跨機箱、氣流因果或 NTC 硬體驗證"
    )
    presentation.save(path)


def sync_pptx_outputs() -> None:
    for path in OUTLINE_PATHS:
        _append_outline(path)
    for path in PPTX_PATHS:
        _append_pptx(path)

